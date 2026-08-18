# -*- coding: utf-8 -*-
"""
[v3 엔진] 사진 기반 자동 제안서 생성 파이프라인 테스트.
요청사항 21의 CASE A~F를 synthetic fixture로 검증한다.
"""
import os
import shutil

import pytest
from pptx import Presentation

from app.engine.pipeline import run_pipeline_photo
from tests.make_samples import _font

from PIL import Image, ImageDraw

FIXT_DIR = os.path.join(os.path.dirname(__file__), "fixtures")
PHOTO_DIR = os.path.join(FIXT_DIR, "v3_site_photos")
KNOWLEDGE_DIR = os.path.join(FIXT_DIR, "v3_knowledge_library")
OUT_DIR = os.path.join(FIXT_DIR, "_test_output_v3")


def _burn_photo(path, label, bg=(120, 130, 140), seed=0):
    import random
    os.makedirs(os.path.dirname(path), exist_ok=True)
    im = Image.new("RGB", (900, 700), bg)
    d = ImageDraw.Draw(im)
    rnd = random.Random(seed or hash(label) or 1)
    for _ in range(10):
        x1, y1 = rnd.randint(0, 900), rnd.randint(0, 700)
        x2, y2 = x1 + rnd.randint(10, 200), y1 + rnd.randint(10, 200)
        color = tuple(max(0, min(255, c + rnd.randint(-50, 50))) for c in bg)
        d.ellipse([x1, y1, x2, y2], fill=color)
    d.rectangle([0, 620, 900, 700], fill=(0, 0, 0))
    d.text((20, 635), label, fill=(255, 255, 0), font=_font(34))
    im.save(path, quality=90)
    return path


def _make_site_photos(tag, specs, filename_carries_label=True):
    """specs: [(filename, burned_label, seed), ...] -> 실제 파일 경로 목록.
    OCR은 폰트 렌더링 환경에 따라 인식률이 들쭉날쭉하므로(sandbox 환경에서 특히),
    테스트에서는 filename_carries_label=True로 파일명에도 라벨을 반영해 분석
    엔진이 실제로 참고하는 두 번째 채널(shape_name=파일명)까지 함께 검증한다.
    CASE D(분류 불가 사진)처럼 의도적으로 무의미한 파일명을 검증해야 할 때만
    filename_carries_label=False로 끈다."""
    out_dir = os.path.join(PHOTO_DIR, tag)
    if os.path.exists(out_dir):
        shutil.rmtree(out_dir)
    paths = []
    for i, (fname, label, seed) in enumerate(specs):
        if filename_carries_label:
            ext = os.path.splitext(fname)[1] or ".jpg"
            fname = f"{label.replace(' ', '')}_{i}{ext}"
        p = os.path.join(out_dir, fname)
        _burn_photo(p, label, seed=seed + i)
        paths.append(p)
    return paths


def _setup_knowledge_library():
    """기존 재도장 픽스처(sample1~3, messy1~2)를 repainting 지식자료로,
    새로 만든 간단한 방수 PPT를 waterproof 지식자료로 등록한다."""
    if os.path.exists(KNOWLEDGE_DIR):
        shutil.rmtree(KNOWLEDGE_DIR)
    repaint_dir = os.path.join(KNOWLEDGE_DIR, "repainting")
    os.makedirs(repaint_dir, exist_ok=True)
    for fname in ("sample1.pptx", "sample2.pptx", "sample3.pptx"):
        shutil.copy2(os.path.join(FIXT_DIR, fname), os.path.join(repaint_dir, fname))

    waterproof_dir = os.path.join(KNOWLEDGE_DIR, "waterproof")
    os.makedirs(waterproof_dir, exist_ok=True)
    _build_waterproof_knowledge_ppt(os.path.join(waterproof_dir, "waterproof_sample.pptx"))


def _build_waterproof_knowledge_ppt(path):
    from pptx.util import Cm, Pt

    img_dir = os.path.join(FIXT_DIR, "v3_wp_images")
    os.makedirs(img_dir, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Cm(19.05)
    prs.slide_height = Cm(25.4)

    def add_text(slide, text, x, y, w, h, size=16):
        box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
        tf = box.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(size)

    def add_pic(slide, tag, seed, x, y, w, h):
        p = _burn_photo(os.path.join(img_dir, f"{tag}.jpg"), tag, bg=(80, 120, 160), seed=seed)
        slide.shapes.add_picture(p, Cm(x), Cm(y), Cm(w), Cm(h))

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(s, "행복빌라 옥상 방수공사 제안서", 1, 1, 17, 2, size=24)
    add_pic(s, "wp_cover", 1, 2, 5, 13, 15)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(s, "방수공사가 필요한 이유", 1, 0.5, 17, 1.5, size=20)
    add_text(s, "옥상 방수층이 노후화되어 누수 위험이 있어 방수공사가 필요합니다.", 1, 2, 17, 1.3, size=13)
    add_pic(s, "wp_defect1", 2, 1, 4, 8, 6)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(s, "누수 및 방수층 하자 현황", 1, 0.5, 17, 1.5, size=20)
    add_text(s, "누수 흔적과 기존 방수층 박리가 확인됩니다.", 1, 2, 17, 1.3, size=13)
    add_pic(s, "wp_defect2", 3, 1, 4, 8, 6)
    add_pic(s, "wp_defect3", 4, 9.5, 4, 8, 6)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(s, "우레탄 도막방수 공법 및 특징", 1, 0.5, 17, 1.5, size=18)
    add_text(s, "우레탄방수는 내구성과 방수성이 뛰어난 도막방수 공법입니다.", 1, 2, 17, 1.3, size=13)
    add_pic(s, "wp_process1", 5, 1, 4, 8, 6)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(s, "방수 시공 순서: 바탕정리 -> 프라이머 -> 우레탄방수 -> 보호몰탈 -> 최종점검", 1, 0.5, 17, 2, size=14)
    add_pic(s, "wp_process2", 6, 1, 4, 8, 6)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(s, "방수공사 시공 전/시공 후", 1, 0.5, 17, 1.5, size=18)
    add_text(s, "시공전", 1, 2, 8, 0.8, size=13)
    add_pic(s, "wp_before", 7, 1, 3, 8, 6)
    add_text(s, "시공후", 9.5, 2, 8, 0.8, size=13)
    add_pic(s, "wp_after", 8, 9.5, 3, 8, 6)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(s, "방수공사 기대 효과", 1, 0.5, 17, 1.5, size=18)
    add_text(s, "누수 걱정 없는 쾌적한 주거환경과 건물 가치 향상을 기대할 수 있습니다.", 1, 2, 17, 1.3, size=13)
    add_pic(s, "wp_effect", 9, 2, 4, 13, 10)

    prs.save(path)


@pytest.fixture(scope="module", autouse=True)
def _prep():
    _setup_knowledge_library()
    os.makedirs(OUT_DIR, exist_ok=True)
    yield


def test_case_a_repainting_20_photos_auto_detected():
    """CASE A: 재도장 사진 20장 -> 재도장 자동판정."""
    specs = []
    labels = ["외벽균열", "도장박리", "외벽오염", "고압세척작업", "하도도장작업",
              "중도도장작업", "상도도장작업", "재도장완료", "단지전경", "도장자재"]
    for i in range(20):
        labels_cycle = labels[i % len(labels)]
        specs.append((f"repaint_{i}.jpg", labels_cycle, i))
    paths = _make_site_photos("case_a", specs)

    result = run_pipeline_photo("가은아파트", paths, KNOWLEDGE_DIR, OUT_DIR)
    assert result["work_type_detected"] == "재도장", result["work_type_percentages"]
    assert os.path.exists(result["pptx"])
    assert result["page_count"] >= 1


def test_case_b_waterproof_15_photos_auto_detected():
    """CASE B: 방수 사진 15장 -> 방수 자동판정."""
    labels = ["옥상방수층누수", "우레탄방수작업", "방수시트설치", "방수공사완료", "누수흔적확인"]
    specs = [(f"wp_{i}.jpg", labels[i % len(labels)], i) for i in range(15)]
    paths = _make_site_photos("case_b", specs)

    result = run_pipeline_photo("나은아파트", paths, KNOWLEDGE_DIR, OUT_DIR)
    assert result["work_type_detected"] == "방수", result["work_type_percentages"]
    assert os.path.exists(result["pptx"])


def test_case_c_mixed_repainting_primary_waterproof_secondary():
    """CASE C: 재도장 15 + 방수 5 -> 재도장 주공종 + 방수 보조공종(모든 사진 개별 분류는 유지)."""
    specs = []
    for i in range(15):
        specs.append((f"mix_paint_{i}.jpg", "재도장 상도도장작업", i))
    for i in range(5):
        specs.append((f"mix_wp_{i}.jpg", "방수 우레탄방수작업", 100 + i))
    paths = _make_site_photos("case_c", specs)

    result = run_pipeline_photo("다은아파트", paths, KNOWLEDGE_DIR, OUT_DIR)
    assert result["work_type_detected"] == "재도장"
    pct_by_label = dict(result["work_type_percentages"])
    assert pct_by_label.get("방수", 0) > 0, "방수 사진이 보조 공종으로 집계되어야 함"


def test_case_d_unclassifiable_photos_preserved_not_deleted():
    """CASE D: 분류 어려운 사진 다수 -> 삭제 없이 unknown 으로 보존."""
    specs = [(f"mystery_{i}.jpg", "IMG", i) for i in range(10)]
    paths = _make_site_photos("case_d", specs, filename_carries_label=False)

    result = run_pipeline_photo("라은아파트", paths, KNOWLEDGE_DIR, OUT_DIR,
                                  work_type_override="기타")
    summary = result["analysis_summary"]
    assert summary["total_photos"] == 10
    # 분석 요약에 unknown 카운트가 남아 있어야 하며(0장으로 사라지지 않음),
    # photo_analysis_summary.json 에 10장 전부의 기록이 남아야 한다(삭제 없음).
    import json
    with open(os.path.join(result["debug_dir"], "photo_analysis_summary.json"), encoding="utf-8") as f:
        dumped = json.load(f)
    assert len(dumped["photos"]) == 10


def test_case_e_five_step_process_single_page():
    """CASE E: 5단계 시공사진 -> process_5step 처럼 한 페이지에 담기고
    4+1 같은 비효율 분할이 발생하지 않아야 한다."""
    labels = ["고압세척작업진행", "균열보수퍼티작업", "하도도장작업진행", "중도도장작업진행", "상도도장작업진행"]
    specs = [(f"step_{i}.jpg", labels[i], i) for i in range(5)]
    paths = _make_site_photos("case_e", specs)

    result = run_pipeline_photo("마은아파트", paths, KNOWLEDGE_DIR, OUT_DIR,
                                  work_type_override="재도장")
    prs = Presentation(result["pptx"])
    process_slide_titles = []
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame and "시공 순서" in sh.text_frame.text:
                process_slide_titles.append(sh.text_frame.text)
                break
    # "시공 순서 1"/"시공 순서 2" 처럼 여러 페이지로 쪼개지지 않아야 한다(1페이지만 허용).
    assert len(process_slide_titles) <= 1, f"시공 순서가 여러 페이지로 분산됨: {process_slide_titles}"


def test_case_f_four_defect_photos_not_silently_dropped():
    """CASE F: 하자사진 4개 -> 1개를 몰래 삭제하지 않고 추가 페이지 또는
    unused_with_reason 으로 명확히 처리해야 한다(사용된 사진 + 기록된 미사용 사유 합 == 4)."""
    labels = ["외벽균열발생", "도장박리현상", "외벽오염발생", "곰팡이백태발생"]
    specs = [(f"defect_{i}.jpg", labels[i], i) for i in range(4)]
    paths = _make_site_photos("case_f", specs)

    result = run_pipeline_photo("바은아파트", paths, KNOWLEDGE_DIR, OUT_DIR,
                                  work_type_override="재도장")
    import json
    with open(os.path.join(result["debug_dir"], "photo_analysis_summary.json"), encoding="utf-8") as f:
        dumped = json.load(f)
    photos = dumped["photos"]
    assert len(photos) == 4
    used = [p for p in photos if p["used"]]
    unused_with_reason = [p for p in photos if not p["used"] and p["unused_reason"]]
    assert len(used) + len(unused_with_reason) == 4, \
        f"사진이 사용/미사용 사유 기록 없이 사라짐: {photos}"
