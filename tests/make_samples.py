# -*- coding: utf-8 -*-
"""
테스트용 가짜 기존 PPT 샘플(2개/3개)을 생성한다.
실제 현장 이름/회사 정보/개인정보가 섞인 상황을 재현하여 파이프라인을 검증하기 위한 픽스처.
"""
import os

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Cm, Pt, Emu

FIXT_DIR = os.path.join(os.path.dirname(__file__), "fixtures")
IMG_DIR = os.path.join(FIXT_DIR, "gen_images")


def _font(size=28):
    try:
        return ImageFont.truetype("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", size, index=2)
    except Exception:
        try:
            return ImageFont.truetype("DejaVuSans.ttf", size)
        except Exception:
            return ImageFont.load_default()


def make_image(path, w, h, bg, label, burn_text=None, badge=None, seed=0):
    import random
    os.makedirs(os.path.dirname(path), exist_ok=True)
    im = Image.new("RGB", (w, h), bg)
    d = ImageDraw.Draw(im)
    rnd = random.Random(seed or hash(label) or 1)
    # 시드 기반 랜덤 도형을 그려 사진마다 지각적 해시가 뚜렷이 달라지도록 함
    for _ in range(14):
        x1, y1 = rnd.randint(0, w), rnd.randint(0, h)
        x2, y2 = x1 + rnd.randint(10, max(11, w // 3)), y1 + rnd.randint(10, max(11, h // 3))
        color = tuple(max(0, min(255, c + rnd.randint(-60, 60))) for c in bg)
        d.ellipse([x1, y1, x2, y2], fill=color)
    d.text((20, 20), label, fill=(255, 255, 255), font=_font(30))
    if burn_text:
        d.rectangle([10, h - 60, w - 10, h - 10], fill=(0, 0, 0))
        d.text((20, h - 50), burn_text, fill=(255, 255, 0), font=_font(22))
    if badge:
        d.rectangle([w - 160, 10, w - 10, 60], fill=(255, 255, 255))
        d.text((w - 150, 18), badge, fill=(200, 30, 30), font=_font(22))
    im.save(path, quality=90)
    return path


def build_images():
    os.makedirs(IMG_DIR, exist_ok=True)
    imgs = {}
    # 하자 사진
    imgs["crack1"] = make_image(f"{IMG_DIR}/crack1.jpg", 900, 700, (120, 110, 100), "외벽 균열 발생 부위", seed=1)
    imgs["crack2"] = make_image(f"{IMG_DIR}/crack2.jpg", 900, 700, (125, 112, 102), "균열 확대 사진", seed=2)
    imgs["peel1"] = make_image(f"{IMG_DIR}/peel1.jpg", 900, 700, (140, 130, 110), "도장 박리 및 들뜸", seed=3)
    imgs["stain1"] = make_image(f"{IMG_DIR}/stain1.jpg", 900, 700, (110, 115, 100), "외벽 오염 및 얼룩", seed=4)
    imgs["mold1"] = make_image(f"{IMG_DIR}/mold1.jpg", 900, 700, (100, 120, 100), "곰팡이 발생 부위", seed=5)

    # 공정 사진
    imgs["wash1"] = make_image(f"{IMG_DIR}/wash1.jpg", 900, 700, (90, 130, 170), "고압세척 작업 중", seed=6)
    imgs["putty1"] = make_image(f"{IMG_DIR}/putty1.jpg", 900, 700, (150, 140, 120), "균열보수 퍼티 작업", seed=7)
    imgs["primer1"] = make_image(f"{IMG_DIR}/primer1.jpg", 900, 700, (170, 170, 175), "하도 도장 작업", seed=8)
    imgs["mid1"] = make_image(f"{IMG_DIR}/mid1.jpg", 900, 700, (180, 180, 60), "중도 도장 작업", seed=9)
    imgs["top1"] = make_image(f"{IMG_DIR}/top1.jpg", 900, 700, (60, 90, 160), "상도 도장 작업", seed=10)
    imgs["safety1"] = make_image(f"{IMG_DIR}/safety1.jpg", 900, 700, (200, 150, 50), "안전 작업 수칙 준수", seed=11)
    imgs["material1"] = make_image(f"{IMG_DIR}/material1.jpg", 900, 700, (210, 210, 200), "도장 자재", seed=12)

    # 전경 / 전후
    imgs["apt_full"] = make_image(
        f"{IMG_DIR}/apt_full.jpg", 1000, 1300, (80, 140, 190), "행복드림아파트 전경",
        burn_text="행복드림아파트 101동", seed=13)
    imgs["before1"] = make_image(
        f"{IMG_DIR}/before1.jpg", 900, 700, (130, 90, 70), "시공 전", burn_text="시공전", seed=14)
    imgs["after1"] = make_image(
        f"{IMG_DIR}/after1.jpg", 900, 700, (70, 150, 210), "시공 후", burn_text="시공후", seed=15)

    # 회사 로고/현수막(제거 대상 테스트)
    imgs["logo"] = make_image(f"{IMG_DIR}/logo.png", 200, 80, (255, 255, 255), "㈜대한페인트건설", seed=16)
    imgs["banner"] = make_image(
        f"{IMG_DIR}/banner.jpg", 1200, 200, (30, 30, 100), "",
        burn_text="㈜대한페인트건설 문의 010-1234-5678", seed=17)

    return imgs


def _add_pic(slide, path, x, y, w, h):
    return slide.shapes.add_picture(path, Cm(x), Cm(y), Cm(w), Cm(h))


def _add_text(slide, text, x, y, w, h, size=18, title=False):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    return box


def build_ppt_1(path):
    prs = Presentation()
    prs.slide_width = Cm(19.05)
    prs.slide_height = Cm(25.4)  # 세로형(4:3 portrait)
    imgs = build_images()

    # 표지: 기존 아파트명/회사명 노출
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, "행복드림아파트 재도장 공사 제안서", 1, 1, 17, 2, size=28, title=True)
    _add_text(s, "서울시 강남구 삼성동 123-45", 1, 3, 17, 1, size=14)
    _add_text(s, "㈜대한페인트건설", 1, 4, 17, 1, size=14)
    _add_pic(s, imgs["apt_full"], 2, 6, 13, 15)
    _add_text(s, "담당자: 홍길동 과장 / 010-1234-5678 / hong@daehanpaint.co.kr", 1, 22, 17, 1.5, size=11)

    # 하자 사진
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, "주요 하자 현황", 1, 0.5, 17, 1.5, size=22, title=True)
    _add_pic(s, imgs["crack1"], 1, 2.5, 8.5, 6.5)
    _add_pic(s, imgs["peel1"], 9.5, 2.5, 8.5, 6.5)
    _add_pic(s, imgs["stain1"], 1, 9.5, 8.5, 6.5)
    _add_pic(s, imgs["mold1"], 9.5, 9.5, 8.5, 6.5)

    # 공법 개요
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, "재도장 공법 개요: 사전점검, 고압세척, 균열보수, 퍼티, 하도, 중도, 상도, 최종점검 순으로 진행", 1, 0.5, 17, 2, size=16, title=True)
    _add_pic(s, imgs["material1"], 1, 3, 8.5, 6.5)
    _add_pic(s, imgs["safety1"], 9.5, 3, 8.5, 6.5)

    # 시공 순서 슬라이드들
    for name, img_key in [("1단계 고압세척", "wash1"), ("2단계 균열보수 퍼티", "putty1"),
                            ("3단계 하도", "primer1"), ("4단계 중도", "mid1"), ("5단계 상도", "top1")]:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _add_text(s, name, 1, 0.5, 17, 1.5, size=20, title=True)
        _add_pic(s, imgs[img_key], 3, 3, 13, 10)

    # 시공 전/후 사례
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, "행복드림아파트 시공 전/후 비교", 1, 0.5, 17, 1.5, size=20, title=True)
    _add_text(s, "시공전", 1, 2, 8, 0.8, size=14)
    _add_pic(s, imgs["before1"], 1, 3, 8.5, 6.5)
    _add_text(s, "시공후", 9.5, 2, 8, 0.8, size=14)
    _add_pic(s, imgs["after1"], 9.5, 3, 8.5, 6.5)

    # 회사 소개 (제거 대상)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_pic(s, imgs["logo"], 1, 1, 6, 2.4)
    _add_text(s, "㈜대한페인트건설 회사 소개", 1, 4, 17, 1.5, size=20, title=True)
    _add_text(s, "홈페이지: www.daehanpaint.co.kr / 사업자등록번호 123-45-67890", 1, 6, 17, 1, size=12)
    _add_pic(s, imgs["banner"], 1, 8, 17, 3)

    prs.save(path)


def build_ppt_2(path):
    prs = Presentation()
    prs.slide_width = Cm(19.05)
    prs.slide_height = Cm(25.4)
    imgs = build_images()

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, "그린빌아파트 외벽 재도장 안내", 1, 1, 17, 2, size=28, title=True)
    _add_text(s, "경기도 성남시 분당구 정자동 88-1 그린빌아파트 관리사무소", 1, 3, 17, 1, size=13)
    _add_pic(s, imgs["apt_full"], 2, 6, 13, 15)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, "하자 및 문제점", 1, 0.5, 17, 1.5, size=22, title=True)
    _add_pic(s, imgs["crack2"], 1, 2.5, 8.5, 6.5)
    _add_pic(s, imgs["stain1"], 9.5, 2.5, 8.5, 6.5)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, "공법 순서: 사전점검 -> 고압세척 -> 균열보수 -> 하도 -> 중도 -> 상도 -> 최종점검", 1, 0.5, 17, 2, size=15, title=True)
    for name, img_key, y in [("고압세척", "wash1", 3), ("균열보수", "putty1", 10)]:
        _add_text(s, name, 1, y, 8, 0.8, size=14)
        _add_pic(s, imgs[img_key], 1, y + 1, 8, 6)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, "하도/중도/상도 도장", 1, 0.5, 17, 1.5, size=20, title=True)
    _add_pic(s, imgs["primer1"], 1, 3, 5.5, 5)
    _add_pic(s, imgs["mid1"], 7, 3, 5.5, 5)
    _add_pic(s, imgs["top1"], 13, 3, 5.5, 5)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, "그린빌아파트 시공 전/시공 후", 1, 0.5, 17, 1.5, size=20, title=True)
    _add_text(s, "Before", 1, 2, 8, 0.8, size=14)
    _add_pic(s, imgs["before1"], 1, 3, 8.5, 6.5)
    _add_text(s, "After", 9.5, 2, 8, 0.8, size=14)
    _add_pic(s, imgs["after1"], 9.5, 3, 8.5, 6.5)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_pic(s, imgs["logo"], 1, 1, 6, 2.4)
    _add_text(s, "시공사: ㈜대한페인트건설 / 담당 이영희 대리 010-9876-5432", 1, 4, 17, 1, size=12)

    prs.save(path)


def build_ppt_3(path):
    """세 번째 파일: 공정 순서가 다르게 표기되고, 사진이 적은 경우 시뮬레이션."""
    prs = Presentation()
    prs.slide_width = Cm(19.05)
    prs.slide_height = Cm(25.4)
    imgs = build_images()

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, "한빛아파트 재도장 공사", 1, 1, 17, 2, size=26, title=True)
    _add_pic(s, imgs["apt_full"], 2, 6, 13, 12)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, "균열 및 박리 하자 사례", 1, 0.5, 17, 1.5, size=20, title=True)
    _add_pic(s, imgs["crack1"], 1, 2.5, 8.5, 6.5)
    _add_pic(s, imgs["peel1"], 9.5, 2.5, 8.5, 6.5)

    # 이 파일은 '중도' 공정 표기가 없음 -> 교차검증 시 needs_confirmation 테스트
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, "공법 순서: 사전점검, 고압세척, 균열보수, 하도, 상도, 최종점검", 1, 0.5, 17, 2, size=15, title=True)
    _add_pic(s, imgs["wash1"], 1, 3, 8, 6)
    _add_pic(s, imgs["top1"], 9.5, 3, 8, 6)

    prs.save(path)


def _generic_photo(path, w, h, bg, tag, seed):
    """실제 회사 PPT에서 흔한, 키워드가 없는 무성의한 사진명(사진1, IMG_001 등)을 재현."""
    return make_image(path, w, h, bg, tag, seed=seed)


def build_messy_ppt(path, variant):
    """사진이 매우 많고(장당 3~4장), 캡션은 대부분 일반적("사진1" 등)이며,
    슬라이드 제목에만 실제 기술 문구가 담긴, 흔한 회사 PPT 스타일을 재현한다.
    회사/현장 정보는 여전히 포함되어 있어 개인정보 제거 로직도 함께 검증한다.
    """
    prs = Presentation()
    prs.slide_width = Cm(19.05)
    prs.slide_height = Cm(25.4)
    img_dir = os.path.join(FIXT_DIR, f"messy_images_{variant}")
    os.makedirs(img_dir, exist_ok=True)

    apt = "무지개마을아파트" if variant == "a" else "코스모스아파트"
    company = "㈜서울외벽방수" if variant == "a" else "동양건업㈜"
    seed_base = 100 if variant == "a" else 300

    def pic(tag, bg, i):
        return _generic_photo(f"{img_dir}/img_{i:03d}.jpg", 900, 700, bg, tag, seed_base + i)

    n = 0
    # 표지(현장정보 포함, 반드시 제거 대상)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, f"{apt} 외벽 재도장 공사", 1, 1, 17, 2, size=26, title=True)
    _add_text(s, f"{company} / 담당 김철수 부장 010-2222-3333", 1, 3, 17, 1, size=12)
    n += 1
    _add_pic(s, pic("현장전경", (90, 130, 175), n), 2, 6, 13, 12)

    # 하자 현황: 캡션은 "사진1"류로 무성의, 제목에만 키워드
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, "외벽 균열 및 오염 등 하자 현황", 1, 0.5, 17, 1.5, size=20, title=True)
    _add_text(s, "외벽 균열부를 V컷팅한 후 탄성퍼티로 보수한다.", 1, 2, 17, 1, size=13)
    positions = [(1, 4, 8.5, 6), (9.5, 4, 8.5, 6), (1, 10.5, 8.5, 6), (9.5, 10.5, 8.5, 6)]
    for i, (x, y, w, h) in enumerate(positions, start=1):
        n += 1
        s2 = pic(f"사진{i}", (120 + i * 5, 110, 100), n)
        _add_pic(s, s2, x, y, w, h)

    # 공법 설명: 실제 기술 문구 + 무성의 캡션 사진 다수
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, "재도장 공법 및 사용 자재 안내", 1, 0.5, 17, 1.5, size=20, title=True)
    _add_text(s, "균열 부위 보수 후 하도·중도·상도 순으로 도장한다. 탄성도료를 사용하여 방수 성능을 함께 확보한다.",
                1, 2, 17, 1.6, size=13)
    for i in range(1, 5):
        n += 1
        x = 1 + (i - 1 % 2) * 0  # simple grid
        row, col = divmod(i - 1, 2)
        s2 = pic(f"IMG_{n:03d}", (150 + i * 3, 140, 120), n)
        _add_pic(s, s2, 1 + col * 8.5, 4 + row * 6.5, 8, 6)

    # 시공 순서 슬라이드 다수 (공정별로 3장씩, 캡션 무성의)
    steps = [("사전 점검", "사전점검을 통해 하자 부위 전체 물량을 산출한다."),
              ("고압세척 작업", "고압수 세정으로 표면 이물질을 완전히 제거한다."),
              ("균열 보수 및 퍼티", "균열부 실링 및 전체 면 퍼티 작업을 진행한다."),
              ("하도 시공", "하도재를 균일하게 도포하여 부착력을 높인다."),
              ("중도 시공", "중도재로 두께를 확보하고 색상을 안정화한다."),
              ("상도 마감", "상도 마감재로 최종 색상과 광택을 완성한다.")]
    for step_name, desc in steps:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _add_text(s, step_name, 1, 0.5, 17, 1.3, size=19, title=True)
        _add_text(s, desc, 1, 2, 17, 1.2, size=12)
        for i in range(1, 4):
            n += 1
            x = 1 + (i - 1) * 6
            s2 = pic(f"현장사진{n}", (100 + n % 50, 120, 140), n)
            _add_pic(s, s2, x, 3.5, 5.5, 5.5)

    # 시공 전/후 사례 2쌍
    for case_i in range(2):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _add_text(s, f"{apt} 시공 전후 비교 {case_i+1}", 1, 0.5, 17, 1.3, size=19, title=True)
        _add_text(s, "시공 전", 1, 2, 8, 0.8, size=13)
        n += 1
        b = pic(f"before_{case_i}", (130, 90, 70 + case_i * 10), n)
        _add_pic(s, b, 1, 3, 8.5, 6.5)
        _add_text(s, "시공 후", 9.5, 2, 8, 0.8, size=13)
        n += 1
        a = pic(f"after_{case_i}", (70, 150, 210 - case_i * 10), n)
        _add_pic(s, a, 9.5, 3, 8.5, 6.5)

    # 완공/전경 추가 사진(효과 페이지용, 대부분 무성의 캡션)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, "시공 완료 후 단지 전경", 1, 0.5, 17, 1.5, size=19, title=True)
    _add_text(s, "재도장 완료 후 외관이 개선되어 입주민 만족도가 향상되었다.", 1, 2, 17, 1.3, size=12)
    for i in range(1, 4):
        n += 1
        s2 = pic(f"완공사진{i}", (80, 150 + i * 5, 190), n)
        _add_pic(s, s2, 1 + (i - 1) * 6, 4, 5.5, 5.5)

    # 회사 소개(제거 대상)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, f"{company} 소개", 1, 1, 17, 1.5, size=20, title=True)
    _add_text(s, f"홈페이지 www.{('seoulwp' if variant=='a' else 'dongyang')}.co.kr "
                  f"사업자등록번호 {'321-45-11111' if variant=='a' else '555-22-33333'}", 1, 3, 17, 1, size=12)

    prs.save(path)
    return n


def main():
    os.makedirs(FIXT_DIR, exist_ok=True)
    build_ppt_1(os.path.join(FIXT_DIR, "sample1.pptx"))
    build_ppt_2(os.path.join(FIXT_DIR, "sample2.pptx"))
    build_ppt_3(os.path.join(FIXT_DIR, "sample3.pptx"))
    n1 = build_messy_ppt(os.path.join(FIXT_DIR, "messy1.pptx"), "a")
    n2 = build_messy_ppt(os.path.join(FIXT_DIR, "messy2.pptx"), "b")
    print("샘플 생성 완료:", FIXT_DIR, f"(messy1={n1}장, messy2={n2}장)")


if __name__ == "__main__":
    main()
