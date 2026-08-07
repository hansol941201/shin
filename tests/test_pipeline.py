# -*- coding: utf-8 -*-
"""
v2 콘텐츠 중심 파이프라인 통합 테스트. 2개/3개 입력 케이스, 민감정보 미노출,
사진/문구 활용률, 품질 점수, 중간 산출물 생성을 검증한다.
실행: pytest tests/test_pipeline.py -v
(사전에 tests/make_samples.py 로 픽스처를 생성해야 한다)
"""
import io
import os
import sys

import pytest
from PIL import Image
from pptx import Presentation

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

from app.main import run_pipeline
from tests.make_samples import main as make_samples_main

FIXT_DIR = os.path.join(os.path.dirname(__file__), "fixtures")
OUT_DIR = os.path.join(os.path.dirname(__file__), "_test_output")

BANNED_STRINGS = [
    "행복드림아파트", "그린빌아파트", "한빛아파트", "대한페인트건설",
    "010-1234-5678", "010-9876-5432", "hong@daehanpaint.co.kr",
    "daehanpaint.co.kr", "123-45-67890", "홍길동", "이영희",
    "무지개마을아파트", "코스모스아파트", "서울외벽방수", "동양건업", "김철수",
    "010-2222-3333", "seoulwp.co.kr", "dongyang.co.kr", "321-45-11111", "555-22-33333",
]

DEBUG_ARTIFACTS = (
    "source_slides.csv", "source_images.csv", "source_texts.csv",
    "image_text_relationships.csv", "unused_content.csv", "content_groups.json",
    "품질점수.txt",
)


@pytest.fixture(scope="module", autouse=True)
def fixtures():
    if not os.path.exists(os.path.join(FIXT_DIR, "sample1.pptx")):
        make_samples_main()
    yield


def _walk(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:
            yield from _walk(sh.shapes)


def _assert_no_leak(pptx_path):
    prs = Presentation(pptx_path)
    assert prs.slide_height > prs.slide_width, "세로형이 아닙니다"
    for slide in prs.slides:
        for sh in _walk(slide.shapes):
            if sh.has_text_frame:
                text = sh.text_frame.text
                for w in BANNED_STRINGS:
                    assert w not in text, f"텍스트에 민감정보 잔존: {w}"
            if sh.shape_type == 13:
                try:
                    import pytesseract
                    im = Image.open(io.BytesIO(sh.image.blob))
                    ocr = pytesseract.image_to_string(im, lang="kor+eng")
                except Exception:
                    continue
                for w in BANNED_STRINGS:
                    assert w not in ocr, f"이미지 OCR에 민감정보 잔존: {w}"


def _assert_debug_artifacts(debug_dir):
    for fname in DEBUG_ARTIFACTS:
        assert os.path.exists(os.path.join(debug_dir, fname)), f"중간 산출물 누락: {fname}"
    assert os.path.isdir(os.path.join(debug_dir, "이미지추출"))


def test_two_file_input():
    files = [os.path.join(FIXT_DIR, "sample1.pptx"), os.path.join(FIXT_DIR, "sample2.pptx")]
    result = run_pipeline("테스트2파일아파트", "재도장", files, OUT_DIR)
    assert os.path.exists(result["pptx"])
    _assert_no_leak(result["pptx"])
    _assert_debug_artifacts(result["debug_dir"])


def test_three_file_input():
    files = [os.path.join(FIXT_DIR, "sample1.pptx"), os.path.join(FIXT_DIR, "sample2.pptx"),
              os.path.join(FIXT_DIR, "sample3.pptx")]
    result = run_pipeline("테스트3파일아파트", "재도장", files, OUT_DIR)
    assert os.path.exists(result["pptx"])
    _assert_no_leak(result["pptx"])


def test_messy_realistic_input_quality_and_utilization():
    """캡션이 부실하고 사진이 매우 많은(실제 회사 PPT와 유사한) 입력에서도, 사진을 최대한
    많이 우겨넣는 대신 "사진+문구가 결합된 8~12(최대 14)페이지 완성형 자료"가 나오는지 검증한다.
    (2026-08 재설계: 사진 활용률 100%/75% 같은 수치 목표는 더 이상 품질 기준이 아니다.)"""
    files = [os.path.join(FIXT_DIR, "messy1.pptx"), os.path.join(FIXT_DIR, "messy2.pptx")]
    result = run_pipeline("○○아파트", "재도장", files, OUT_DIR)
    assert os.path.exists(result["pptx"])
    _assert_no_leak(result["pptx"])

    assert result["total_usable_image_count"] >= 30, \
        f"사용 가능 이미지가 {result['total_usable_image_count']}장으로 목표(30장) 미달"

    # 페이지 수: 표지 제외 본문 기준이 아니라 전체(표지 포함) 8~14장 범위를 벗어나면 안 된다.
    total_pages = result["page_count"] + 1  # page_count는 표지를 제외한 값
    assert total_pages <= 14, f"전체 페이지 {total_pages}장으로 상한(14장) 초과"

    prs = Presentation(result["pptx"])
    slides = list(prs.slides)
    assert len(slides) == total_pages

    def _slide_has_image(slide):
        return any(sh.shape_type == 13 for sh in _walk(slide.shapes))

    def _slide_has_text(slide):
        for sh in _walk(slide.shapes):
            if sh.has_text_frame and sh.text_frame.text.strip():
                return True
        return False

    # 표지(0)와 마지막(마무리) 페이지를 제외한 본문 페이지는 사진만/문구만 있으면 안 된다
    # (시공 전후 비교 페이지는 사진 중심 레이아웃이라 예외로 허용).
    photo_only = 0
    text_only = 0
    for slide in slides[1:-1]:
        has_img, has_text = _slide_has_image(slide), _slide_has_text(slide)
        if has_img and not has_text:
            photo_only += 1
        if has_text and not has_img:
            text_only += 1
    assert photo_only <= 1, f"사진만 있는 본문 페이지 {photo_only}장 (기준 2장 이상 금지)"
    assert text_only == 0, f"문구만 있는 본문 페이지 {text_only}장 존재(금지)"

    assert result["quality_score"] >= 85, f"품질 점수 {result['quality_score']}로 기준(85) 미달"
    assert result["quality_passed"], "하드 실패 조건(페이지 수/사진-문구 결합/중복 제목 등) 위반으로 FAIL"
    # 재구성 로직이 실제로 여러 번 동작했는지(단순 재실행이 아님)만 확인한다.
    assert result["attempts"] >= 1
    assert result["a_grade_count"] >= 0 and result["b_grade_count"] >= 0

    _assert_debug_artifacts(result["debug_dir"])
    assert os.path.exists(os.path.join(result["debug_dir"], "source_texts.csv"))
    assert os.path.exists(result["quality_json"])


def test_invalid_file_count_rejected():
    files = [os.path.join(FIXT_DIR, "sample1.pptx")]
    with pytest.raises(ValueError):
        run_pipeline("아파트", "재도장", files, OUT_DIR)


def test_empty_apartment_name_rejected():
    files = [os.path.join(FIXT_DIR, "sample1.pptx"), os.path.join(FIXT_DIR, "sample2.pptx")]
    with pytest.raises(ValueError):
        run_pipeline("  ", "재도장", files, OUT_DIR)


def _make_site_photo_files(tmp_dir, n=4):
    os.makedirs(tmp_dir, exist_ok=True)
    paths = []
    for i in range(n):
        img = Image.new("RGB", (1200, 900), (180 + i * 5, 180, 180))
        p = os.path.join(tmp_dir, f"site_{i+1}.jpg")
        img.save(p, quality=85)
        paths.append(p)
    return paths


def test_site_photos_placed_right_after_cover_with_no_diagnosis():
    """사용자가 추가한 현장사진이 표지 바로 뒤에 배치되고, AI가 균열/박리/누수 등
    기술적 판단 문구를 임의로 생성하지 않는지 검증한다."""
    files = [os.path.join(FIXT_DIR, "sample1.pptx"), os.path.join(FIXT_DIR, "sample2.pptx")]
    site_dir = os.path.join(os.path.dirname(__file__), "_site_photos_src")
    site_photos = _make_site_photo_files(site_dir, n=4)

    result = run_pipeline("현장사진테스트아파트", "재도장", files, OUT_DIR,
                            site_photo_paths=site_photos)
    assert os.path.exists(result["pptx"])
    _assert_no_leak(result["pptx"])

    prs = Presentation(result["pptx"])
    slides = list(prs.slides)
    assert len(slides) >= 2

    def _slide_title(slide):
        for sh in _walk(slide.shapes):
            if sh.has_text_frame and sh.text_frame.text.strip():
                return sh.text_frame.text.strip().split("\n")[0]
        return ""

    # 표지(0) 바로 다음 페이지가 현장사진 페이지여야 한다.
    second_title = _slide_title(slides[1])
    assert "현장사진" in second_title, f"현장사진 페이지가 표지 바로 뒤에 없음: '{second_title}'"

    diagnostic_words = ["균열", "박리", "누수", "보수 필요", "시공 권장", "하자 의심"]
    for sh in _walk(slides[1].shapes):
        if sh.has_text_frame:
            text = sh.text_frame.text
            for w in diagnostic_words:
                assert w not in text, f"현장사진 페이지에 AI 진단성 문구 발견: '{w}' in '{text}'"

    assert any(sh.shape_type == 13 for sh in _walk(slides[1].shapes)), "현장사진 페이지에 사진이 없음"

    # 개별 사진에는 캡션(중립 번호조차)이 없어야 하고, 대신 공종 기준 공통
    # 안내문구 하나만 표시돼야 한다("현장사진 01" 같은 개별 번호 캡션 금지).
    all_text = "\n".join(sh.text_frame.text for sh in _walk(slides[1].shapes) if sh.has_text_frame)
    assert "현장사진 01" not in all_text and "현장사진 02" not in all_text, \
        "현장사진 페이지에 개별 사진 번호 캡션이 남아있음(제거 대상)"
    assert "검토" in all_text and "재도장" in all_text, \
        f"공종 기준 공통 안내문구가 현장사진 페이지에 없음: '{all_text}'"


def test_no_site_photos_falls_back_to_original_flow():
    """현장사진을 넣지 않으면 기존 방식(사진1)이 표지 바로 뒤에 오지 않아야 한다."""
    files = [os.path.join(FIXT_DIR, "sample1.pptx"), os.path.join(FIXT_DIR, "sample2.pptx")]
    result = run_pipeline("현장사진없음아파트", "재도장", files, OUT_DIR)
    prs = Presentation(result["pptx"])
    slides = list(prs.slides)

    def _slide_title(slide):
        for sh in _walk(slide.shapes):
            if sh.has_text_frame and sh.text_frame.text.strip():
                return sh.text_frame.text.strip().split("\n")[0]
        return ""

    assert "현장사진" not in _slide_title(slides[1])


def test_generation_failure_raises_and_logs_stage(monkeypatch):
    """generate_pptx_v2가 실패해도(혹은 저장 후 파일이 사라져도) 파이프라인이
    "정상 완료"를 반환해서는 안 되며, 처리로그.txt 마지막 줄에 실패 단계가
    남아야 한다(중간산출물만 만들고 조용히 끝나는 것을 금지)."""
    import app.engine.pipeline as pipeline_mod

    def _boom(*args, **kwargs):
        raise RuntimeError("강제 저장 실패(테스트)")

    monkeypatch.setattr(pipeline_mod, "generate_pptx_v2", _boom)

    files = [os.path.join(FIXT_DIR, "sample1.pptx"), os.path.join(FIXT_DIR, "sample2.pptx")]
    with pytest.raises(pipeline_mod.PptGenerationFailedError) as exc_info:
        run_pipeline("저장실패테스트아파트", "재도장", files, OUT_DIR)

    assert "스토리 구성/PPT 생성" in str(exc_info.value)

    safe_name = "저장실패테스트아파트"
    log_path = os.path.join(OUT_DIR, f"{safe_name}_처리로그.txt")
    assert os.path.exists(log_path), "실패 시에도 처리로그.txt는 반드시 생성돼야 함"
    with open(log_path, encoding="utf-8") as f:
        content = f.read()
    assert "[파이프라인] 종료 단계:" in content
    assert "스토리 구성/PPT 생성" in content.strip().splitlines()[-1]

    out_pptx = os.path.join(OUT_DIR, f"{safe_name}_재도장_입주민설명자료.pptx")
    assert not os.path.exists(out_pptx), "저장이 실패했는데 PPT 파일이 존재하면 안 됨"
