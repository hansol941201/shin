# -*- coding: utf-8 -*-
"""
결과물 저장 전 자동 검수를. 문제 발견 시 가능한 범위에서 자동 수정하고,
자동 수정이 불가능한 항목은 검수 보고서에 기록한다.
"""
import io
from dataclasses import dataclass, field
from typing import List

import imagehash
from PIL import Image
from pptx import Presentation

from app.anonymizer.anonymizer import Blacklist, text_contains_banned
from app.utils.config import RE_EMAIL, RE_PHONE, RE_URL, RE_BIZNO


@dataclass
class ValidationReport:
    passed: List[str] = field(default_factory=list)
    autofixed: List[str] = field(default_factory=list)
    manual_review: List[str] = field(default_factory=list)

    def as_text(self) -> str:
        lines = ["=== 검수 결과 보고서 ===", ""]
        lines.append(f"[통과] {len(self.passed)}건")
        for p in self.passed:
            lines.append(f"  - {p}")
        lines.append("")
        lines.append(f"[자동 수정] {len(self.autofixed)}건")
        for p in self.autofixed:
            lines.append(f"  - {p}")
        lines.append("")
        lines.append(f"[수동 확인 필요] {len(self.manual_review)}건")
        for p in self.manual_review:
            lines.append(f"  - {p}")
        return "\n".join(lines)


def _rect_overlap_ratio(a, b) -> float:
    """두 사각형(left, top, width, height, EMU)이 겹치는 비율(작은 쪽 면적 대비)을 반환한다."""
    ax0, ay0, ax1, ay1 = a[0], a[1], a[0] + a[2], a[1] + a[3]
    bx0, by0, bx1, by1 = b[0], b[1], b[0] + b[2], b[1] + b[3]
    iw = max(0, min(ax1, bx1) - max(ax0, bx0))
    ih = max(0, min(ay1, by1) - max(ay0, by0))
    inter = iw * ih
    if inter <= 0:
        return 0.0
    area_a, area_b = a[2] * a[3], b[2] * b[3]
    smaller = min(area_a, area_b)
    return inter / smaller if smaller > 0 else 0.0


def _check_text_overlaps(prs) -> int:
    """슬라이드별로 텍스트가 있는 도형끼리 겹치는지 검사한다(레이아웃 계산
    실수로 두 텍스트 상자가 서로 겹치는 실제 버그를 잡기 위함). 겹침 비율이
    25%를 넘는 쌍을 카운트한다."""
    overlap_count = 0
    for slide in prs.slides:
        boxes = []
        for shape in _iter_all_shapes(slide.shapes):
            try:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
                        continue
                    boxes.append((shape.left, shape.top, shape.width, shape.height))
            except Exception:
                pass
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                if _rect_overlap_ratio(boxes[i], boxes[j]) > 0.25:
                    overlap_count += 1
    return overlap_count


def _iter_all_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:
            try:
                yield from _iter_all_shapes(shape.shapes)
            except Exception:
                pass


def validate_and_fix(pptx_path: str, blacklist: Blacklist, new_apartment_name: str) -> ValidationReport:
    report = ValidationReport()
    prs = Presentation(pptx_path)
    modified = False

    # 1) 세로형 검사
    if prs.slide_height >= prs.slide_width:
        report.passed.append("모든 슬라이드가 세로형(A4 비율)입니다.")
    else:
        report.manual_review.append("슬라이드가 가로형으로 생성되었습니다.")

    apt_name_found_on_cover = False
    leaked_terms = set()
    off_slide_count = 0
    seen_hashes = []
    duplicate_images = 0

    for s_idx, slide in enumerate(prs.slides):
        for shape in _iter_all_shapes(slide.shapes):
            # 텍스트 검사
            try:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if text:
                        if new_apartment_name and new_apartment_name in text:
                            apt_name_found_on_cover = True
                        for label, pattern in (("전화번호", RE_PHONE), ("이메일", RE_EMAIL),
                                                ("홈페이지 주소", RE_URL), ("사업자등록번호", RE_BIZNO)):
                            if pattern.search(text):
                                leaked_terms.add(label)
                        for name in blacklist.apt_names:
                            if name and name in text:
                                leaked_terms.add(f"기존 아파트명({name})")
                        for name in blacklist.companies:
                            if name and name in text:
                                leaked_terms.add(f"회사명({name})")
            except Exception:
                pass

            # 슬라이드 밖으로 나간 요소 검사 + 자동 클램프
            try:
                left, top = shape.left, shape.top
                width, height = shape.width, shape.height
                if left is None or top is None or width is None or height is None:
                    continue
                out_of_bounds = (left < 0 or top < 0 or
                                  left + width > prs.slide_width + 1000 or
                                  top + height > prs.slide_height + 1000)
                if out_of_bounds:
                    off_slide_count += 1
                    new_left = max(0, min(int(left), int(prs.slide_width - width)))
                    new_top = max(0, min(int(top), int(prs.slide_height - height)))
                    shape.left, shape.top = new_left, new_top
                    modified = True
            except Exception:
                pass

            # 이미지 중복/해상도 검사
            try:
                if shape.shape_type == 13:  # PICTURE
                    blob = shape.image.blob
                    im = Image.open(io.BytesIO(blob)).convert("RGB")
                    h = imagehash.phash(im)
                    dup = any((h - existing) <= 3 for existing in seen_hashes)
                    if dup:
                        duplicate_images += 1
                    else:
                        seen_hashes.append(h)

                    # 과도한 확대 검사: 표시 크기(px 환산)가 원본 대비 1.6배 초과
                    disp_w_px = shape.width / 914400 * 96
                    disp_h_px = shape.height / 914400 * 96
                    if disp_w_px > im.width * 1.6 or disp_h_px > im.height * 1.6:
                        report.manual_review.append(
                            f"{s_idx+1}번 슬라이드: 저해상도 이미지가 과도하게 확대되었을 수 있습니다.")
            except Exception:
                pass

    if new_apartment_name:
        if apt_name_found_on_cover:
            report.passed.append("새 아파트명이 문서 내에 정상적으로 적용되었습니다.")
        else:
            report.manual_review.append("새 아파트명이 슬라이드 텍스트에서 확인되지 않았습니다.")

    if leaked_terms:
        report.manual_review.append(f"제거되지 않은 민감 정보 발견: {', '.join(sorted(leaked_terms))}")
    else:
        report.passed.append("기존 아파트명·회사명·연락처 등 민감 정보가 발견되지 않았습니다.")

    if off_slide_count:
        report.autofixed.append(f"슬라이드 밖으로 벗어난 요소 {off_slide_count}건을 슬라이드 안쪽으로 재배치했습니다.")
    else:
        report.passed.append("모든 요소가 슬라이드 영역 안에 위치합니다.")

    if duplicate_images:
        report.manual_review.append(f"결과물 내 반복되는 사진이 {duplicate_images}건 발견되었습니다.")
    else:
        report.passed.append("동일하거나 거의 유사한 사진의 중복 사용이 발견되지 않았습니다.")

    # 마지막 슬라이드에 회사 정보가 없는지 별도 재확인
    try:
        last = list(prs.slides)[-1]
        last_text = ""
        for shape in _iter_all_shapes(last.shapes):
            if shape.has_text_frame:
                last_text += shape.text_frame.text
        if text_contains_banned(last_text, blacklist) or RE_PHONE.search(last_text) or RE_EMAIL.search(last_text):
            report.manual_review.append("마지막 페이지에 회사/연락처 정보로 의심되는 텍스트가 남아 있습니다.")
        else:
            report.passed.append("마지막 페이지에 회사 정보가 포함되어 있지 않습니다.")
    except Exception:
        pass

    # 텍스트 상자 겹침 검사(레이아웃 계산 버그로 두 텍스트가 서로 겹쳐 보이는지)
    overlap_count = _check_text_overlaps(prs)
    if overlap_count:
        report.manual_review.append(
            f"텍스트 겹침 의심 {overlap_count}건이 발견되었습니다 - 페이지 레이아웃을 확인해주세요.")
    else:
        report.passed.append("텍스트 상자 겹침이 발견되지 않았습니다.")

    if modified:
        prs.save(pptx_path)

    return report
