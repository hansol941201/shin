# -*- coding: utf-8 -*-
"""v2 엔진 전용 저수준 렌더링 유틸리티(legacy와 독립적으로 유지)."""
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Mm, Pt

from app.utils.config import (
    SLIDE_WIDTH, SLIDE_HEIGHT, COLOR_WHITE, COLOR_NAVY, COLOR_NAVY_LIGHT,
    COLOR_GRAY_LIGHT, COLOR_GRAY_MID, COLOR_GOLD, COLOR_TEXT_BODY, COLOR_TEXT_SUB,
    FONT_PRIMARY, FONT_SIZE_PAGE_TITLE, FONT_SIZE_BODY, FONT_SIZE_CAPTION, FONT_SIZE_FOOTER,
)

MARGIN = Mm(12)
CONTENT_W = SLIDE_WIDTH - MARGIN * 2


def _rgb(t):
    return RGBColor(*t)


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(COLOR_WHITE)
    return slide


def _set_font(run, size=FONT_SIZE_BODY, color=COLOR_TEXT_BODY, bold=False, name=FONT_PRIMARY):
    run.font.size = size
    run.font.bold = bold
    run.font.name = name
    run.font.color.rgb = _rgb(color)


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, size=FONT_SIZE_BODY, color=COLOR_TEXT_BODY,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        _set_font(run, size=size, color=color, bold=bold)
    return box


def _fit_box(img_w, img_h, max_w, max_h):
    ratio = min(max_w / img_w, max_h / img_h)
    return int(img_w * ratio), int(img_h * ratio)


def add_picture_contain(slide, path, left, top, max_w, max_h, frame=True):
    """비율을 유지한 채 지정 영역 안에 이미지를 중앙 배치(찌그러뜨리지 않음)."""
    try:
        with Image.open(path) as im:
            iw, ih = im.size
    except Exception:
        iw, ih = max_w, max_h
    w, h = _fit_box(iw, ih, max_w, max_h)
    x = left + (max_w - w) // 2
    y = top + (max_h - h) // 2
    if frame:
        pad = Emu(9525 * 2)
        add_rect(slide, x - pad, y - pad, w + pad * 2, h + pad * 2, COLOR_GRAY_LIGHT)
    pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
    return pic, (x, y, w, h)


def header(slide, title_text, page_no=None):
    add_rect(slide, 0, 0, SLIDE_WIDTH, Mm(20), COLOR_NAVY)
    add_rect(slide, 0, Mm(20), SLIDE_WIDTH, Pt(2.2), COLOR_GOLD)
    add_text(slide, MARGIN, Mm(3.5), CONTENT_W, Mm(13), title_text,
              size=FONT_SIZE_PAGE_TITLE, color=COLOR_WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if page_no is not None:
        add_text(slide, SLIDE_WIDTH - Mm(30), SLIDE_HEIGHT - Mm(12), Mm(20), Mm(8),
                  f"{page_no:02d}", size=FONT_SIZE_FOOTER, color=COLOR_GRAY_MID, align=PP_ALIGN.RIGHT)
