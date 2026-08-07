# -*- coding: utf-8 -*-
"""
Placeholder Engine: 디자이너가 만든 PowerPoint 템플릿(.pptx) 파일을 열어,
그 안의 {{PLACEHOLDER}} 텍스트와 PHOTO_N 이름의 도형만 실제 값으로 교체한다.

절대 하지 않는 것:
- 도형을 새로 그리거나 색상/폰트/여백/정렬을 바꾸는 것(디자인은 템플릿이 담당)
- 사진을 찌그러뜨리거나 원본 비율을 벗어나게 확대하는 것

값이 제공되지 않은 선택적 placeholder(예: 카드 4개짜리 템플릿인데 문구가
3개뿐인 경우)는 해당 카드 그룹 전체(사진+제목+설명)를 통째로 제거해
빈 자리가 남지 않도록 한다.
"""
import copy
import re
from typing import Dict, Optional

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

TOKEN_RE = re.compile(r"^\{\{([A-Z0-9_]+)\}\}$")


def _fit_box(img_w, img_h, max_w, max_h):
    ratio = min(max_w / img_w, max_h / img_h)
    return int(img_w * ratio), int(img_h * ratio)


def _insert_picture(shapes, image_path: str, left, top, width, height):
    """비율을 유지한 채(찌그러뜨리지 않고, 과도하게 확대하지 않고) 지정 영역
    안에 사진을 중앙 배치한다."""
    try:
        with Image.open(image_path) as im:
            iw, ih = im.size
    except Exception:
        iw, ih = width, height
    w, h = _fit_box(iw, ih, width, height)
    x = left + (width - w) // 2
    y = top + (height - h) // 2
    return shapes.add_picture(image_path, int(x), int(y), width=int(w), height=int(h))


def _shape_token_values(shape, text_map: Dict[str, str]):
    """도형 안의 모든 placeholder 토큰과, text_map에 실제 채울 값이 있는지 여부를 반환."""
    tokens = []
    if not shape.has_text_frame:
        return tokens
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            m = TOKEN_RE.match(run.text.strip())
            if m:
                tokens.append((run, m.group(1)))
    return tokens


def _fill_text_shapes(shape_iterable, text_map: Dict[str, str]):
    for shape in list(shape_iterable):
        for run, key in _shape_token_values(shape, text_map):
            run.text = text_map.get(key) or ""


def _photo_name(shape) -> Optional[str]:
    name = shape.name or ""
    return name if name.startswith("PHOTO_") else None


def _fill_photo_shapes(container_shapes, photo_map: Dict[str, str]):
    """PHOTO_N 이름의 사각형 자리표시자를 실제 사진으로 교체한다(같은 위치/
    크기에 비율 유지 삽입 후 원래 자리표시자 도형은 제거)."""
    for shape in list(container_shapes):
        pname = _photo_name(shape)
        if not pname:
            continue
        path = photo_map.get(pname)
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        parent = shape._element.getparent()
        idx = list(parent).index(shape._element)
        parent.remove(shape._element)
        if path:
            pic = _insert_picture(container_shapes, path, left, top, width, height)
            # 원래 자리표시자가 있던 z-order 위치로 옮겨 다른 도형과의 겹침 순서를 보존한다.
            parent.remove(pic._element)
            parent.insert(idx, pic._element)


def _group_has_value(group_shape, text_map: Dict[str, str], photo_map: Dict[str, str]) -> bool:
    for sh in group_shape.shapes:
        for _run, key in _shape_token_values(sh, text_map):
            if text_map.get(key):
                return True
        pname = _photo_name(sh)
        if pname and photo_map.get(pname):
            return True
    return False


def fill_template(template_path: str, text_map: Dict[str, str], photo_map: Dict[str, str]):
    """템플릿을 열어 placeholder를 채운 Presentation 객체를 반환한다(템플릿
    파일 자체는 수정하지 않음 - 매번 새로 읽어서 채운다)."""
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # 1) 그룹(카드/스텝 등 반복 단위) 처리: 값이 하나도 없으면 그룹 전체 제거,
    #    값이 있으면 그룹 내부의 텍스트/사진 placeholder만 채운다.
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            if not _group_has_value(shape, text_map, photo_map):
                shape._element.getparent().remove(shape._element)
            else:
                _fill_text_shapes(shape.shapes, text_map)
                _fill_photo_shapes(shape.shapes, photo_map)

    # 2) 그룹에 속하지 않은 단독 텍스트/사진 placeholder 처리
    _fill_text_shapes(slide.shapes, text_map)
    _fill_photo_shapes(slide.shapes, photo_map)

    return prs
