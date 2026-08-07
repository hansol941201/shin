# -*- coding: utf-8 -*-
"""
서로 다른 Presentation(템플릿 파일 vs 최종 출력 파일) 사이에서 슬라이드를
그대로 복제하는 유틸리티. python-pptx는 프레젠테이션 간 슬라이드 복사를
기본 지원하지 않으므로, 도형 XML을 통째로 복사한 뒤 그 안에서 참조하는
이미지 파트(rId)만 대상 프레젠테이션에 다시 등록(rehome)해 연결을 살린다.

이렇게 해야 templates/*.pptx(디자이너가 만든 실제 PowerPoint 파일)의 모든
도형·서식·색상·SmartArt 등을 원본 그대로 최종 결과물에 옮길 수 있다 - 도형을
파이썬 코드로 다시 그리는 것이 아니라, 템플릿 파일 자체를 재사용한다.
"""
import copy
from typing import Dict

from pptx.oxml.ns import qn

_BLIP_TAGS = (qn("a:blip"),)
_EMBED_ATTR = qn("r:embed")
_LINK_ATTR = qn("r:link")


def _rehome_images(dest_slide, src_slide, element) -> None:
    """복사된 XML 조각(element) 안에서 이미지(blip) 참조를 찾아, 원본
    슬라이드의 이미지 데이터를 대상 슬라이드 파트에 새로 등록하고 rId를
    갱신한다(그래야 복사된 도형이 올바른 이미지를 계속 가리킨다)."""
    rid_map: Dict[str, str] = {}
    for blip in element.iter(qn("a:blip")):
        old_rid = blip.get(_EMBED_ATTR)
        if not old_rid:
            continue
        if old_rid in rid_map:
            blip.set(_EMBED_ATTR, rid_map[old_rid])
            continue
        try:
            image_part = src_slide.part.related_part(old_rid)
            new_image_part, new_rid = dest_slide.part.get_or_add_image_part(
                _ImageStreamAdapter(image_part))
        except Exception:
            continue
        rid_map[old_rid] = new_rid
        blip.set(_EMBED_ATTR, new_rid)


class _ImageStreamAdapter:
    """get_or_add_image_part는 파일 경로 또는 파일 객체를 받는다. 이미 로드된
    ImagePart의 바이트를 파일처럼 넘겨주기 위한 얇은 어댑터."""
    def __init__(self, image_part):
        import io
        self._stream = io.BytesIO(image_part.blob)

    def read(self, *a, **kw):
        return self._stream.read(*a, **kw)

    def seek(self, *a, **kw):
        return self._stream.seek(*a, **kw)

    def tell(self, *a, **kw):
        return self._stream.tell(*a, **kw)


def clone_slide_into(dest_prs, src_slide):
    """src_slide(다른 Presentation 객체에 속함)의 모든 도형을 dest_prs에
    새로 추가한 빈 슬라이드로 복제한다. 배경색/도형/텍스트 서식/이미지가
    모두 그대로 유지된다. 새로 생성된 dest 슬라이드를 반환한다."""
    blank_layout = dest_prs.slide_layouts[6]
    dest_slide = dest_prs.slides.add_slide(blank_layout)

    # 배경(슬라이드 고유 배경이 있으면 함께 복제)
    try:
        src_bg = src_slide.element.find(qn("p:bg"))
        if src_bg is not None:
            dest_slide.element.insert(0, copy.deepcopy(src_bg))
    except Exception:
        pass

    for shape in list(src_slide.shapes):
        new_el = copy.deepcopy(shape._element)
        _rehome_images(dest_slide, src_slide, new_el)
        dest_slide.shapes._spTree.append(new_el)

    return dest_slide
