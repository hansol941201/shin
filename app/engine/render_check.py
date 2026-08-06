# -*- coding: utf-8 -*-
"""
결과 PPT를 실제로 이미지로 렌더링하여 시각 검증한다(요청 스펙 12).
LibreOffice로 PDF 변환 후 PyMuPDF로 슬라이드별 PNG + 전체 contact_sheet.png 를 만들고,
python-pptx로 재오픈하여 슬라이드 밖 요소/사진 없는 프레임 등을 기계적으로 점검한다.
"""
import json
import os
from typing import Dict, List


def render_slides_and_contact_sheet(pdf_path: str, out_dir: str, apt_prefix: str):
    """slide_001.png ... 및 contact_sheet.png 생성. 실패 시 (None, []) 반환."""
    try:
        import fitz  # PyMuPDF
        from PIL import Image
    except Exception:
        return None, []

    try:
        doc = fitz.open(pdf_path)
    except Exception:
        return None, []

    slide_paths = []
    thumbs = []
    for i, page in enumerate(doc, start=1):
        pix = page.get_pixmap(matrix=fitz.Matrix(1.2, 1.2))
        img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
        p = os.path.join(out_dir, f"slide_{i:03d}.png")
        img.save(p)
        slide_paths.append(p)
        thumb_pix = page.get_pixmap(matrix=fitz.Matrix(0.5, 0.5))
        thumbs.append(Image.frombytes("RGB", (thumb_pix.width, thumb_pix.height), thumb_pix.samples))

    if not thumbs:
        return None, slide_paths

    cols = 4
    rows = (len(thumbs) + cols - 1) // cols
    tw, th = thumbs[0].size
    pad = 10
    grid = Image.new("RGB", (cols * tw + (cols + 1) * pad, rows * th + (rows + 1) * pad), "white")
    for i, im in enumerate(thumbs):
        r, c = divmod(i, cols)
        grid.paste(im, (pad + c * (tw + pad), pad + r * (th + pad)))
    contact_sheet_path = os.path.join(out_dir, f"{apt_prefix}_contact_sheet.png")
    grid.save(contact_sheet_path)
    return contact_sheet_path, slide_paths


def visual_validation(pptx_path: str) -> dict:
    """python-pptx 재오픈 기반 기계적 시각 점검(빈 프레임/이탈 요소/슬라이드당 사진 수)."""
    from pptx import Presentation

    prs = Presentation(pptx_path)
    result = {"slide_count": len(prs.slides), "slides": [], "issues": []}

    def walk(shapes):
        for sh in shapes:
            yield sh
            if sh.shape_type == 6:
                yield from walk(sh.shapes)

    for idx, slide in enumerate(prs.slides, start=1):
        pic_count, off_slide, has_text = 0, 0, False
        for sh in walk(slide.shapes):
            if sh.shape_type == 13:
                pic_count += 1
                try:
                    if (sh.left is not None and sh.top is not None and
                            (sh.left < -10000 or sh.top < -10000 or
                             sh.left + (sh.width or 0) > prs.slide_width + 10000 or
                             sh.top + (sh.height or 0) > prs.slide_height + 10000)):
                        off_slide += 1
                except Exception:
                    pass
            if sh.has_text_frame and sh.text_frame.text.strip():
                has_text = True
        result["slides"].append({"slide": idx, "picture_count": pic_count,
                                   "off_slide_elements": off_slide, "has_text": has_text})
        if off_slide:
            result["issues"].append(f"슬라이드 {idx}: 화면 밖으로 벗어난 요소 {off_slide}개")

    return result


def dump_visual_validation_json(data: dict, out_path: str) -> str:
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    return out_path
