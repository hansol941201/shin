# -*- coding: utf-8 -*-
"""
v2 콘텐츠 중심 PPTX 렌더러.

[재설계 2 - 2026-08] 모든 페이지를 같은 "사진 그리드 + 문구 박스" 틀로 찍어내지 않는다.
story.py가 각 페이지에 지정한 semantic_type(hero/reason_hero/four_cards/
image_text_split/feature_cards/process_timeline/before_after/two_case_compare/
effects_hero/closing)에 따라 서로 다른 전용 레이아웃을 그린다.

[디자인 시스템 - 2026-08] semantic_type 하나당 여러 개의 시각적 "레이아웃 변형(variant)"을
LAYOUT_CATALOG에 등록해 두고(총 20개 이상), generate_pptx_v2가 페이지를 순서대로
그리면서 직전 페이지와 같은 변형이 연속으로 나오지 않도록 자동으로 다른 변형을
고른다. 카드형/타임라인/Before-After/숫자 강조/매거진 스타일 등을 섞어 "AI가 찍어낸
PPT"가 아니라 사람이 페이지마다 다르게 편집한 것처럼 보이게 하는 것이 목표다.
"""
import os
import re
from typing import Dict, List

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Mm, Pt

from app.engine.render_utils import (
    MARGIN, CONTENT_W, new_presentation, blank_slide, add_rect, add_text,
    add_picture_contain, header,
)
from app.utils.config import (
    SLIDE_WIDTH, SLIDE_HEIGHT, COLOR_WHITE, COLOR_NAVY, COLOR_NAVY_LIGHT, COLOR_GRAY_LIGHT,
    COLOR_GRAY_MID, COLOR_GOLD, COLOR_TEXT_BODY, COLOR_TEXT_SUB, FONT_SIZE_CAPTION,
    FONT_SIZE_BODY, FONT_SIZE_SECTION_HEAD,
)

MM = Mm(1)


# ------------------------------------------------------------------
# 공통 저수준 조각(여러 semantic_type 레이아웃에서 재사용)
# ------------------------------------------------------------------

def _grid_geometry(n: int):
    if n <= 1:
        return 1, 1
    if n == 2:
        return 2, 1
    if n == 3:
        return 3, 1
    if n == 4:
        return 2, 2
    return 3, (n + 2) // 3


def _photo_grid(slide, images: List, top, area_h, cols=None, rows=None, show_caption=True):
    n = len(images)
    if n == 0:
        return
    if cols is None or rows is None:
        cols, rows = _grid_geometry(n)
    gap = Mm(4)
    cap_h = Mm(9) if show_caption else Mm(0)
    cell_w = (CONTENT_W - gap * (cols - 1)) / cols
    cell_h = (area_h - gap * (rows - 1)) / rows
    for idx, img in enumerate(images):
        r, c = divmod(idx, cols)
        x = MARGIN + Emu(int(cell_w + gap) * c)
        y = top + Emu(int(cell_h + gap) * r)
        add_picture_contain(slide, img.path, x, y, int(cell_w), int(cell_h) - cap_h)
        if show_caption:
            cap = img.real_caption or "현장 시공 상태 참고사진"
            add_text(slide, x, y + int(cell_h) - cap_h, int(cell_w), cap_h, cap,
                      size=FONT_SIZE_CAPTION, color=COLOR_TEXT_SUB, align=PP_ALIGN.CENTER)


def _photo_row(slide, images: List, top, height, show_caption=True):
    """가로 한 줄(썸네일 행). support_images 등 보조 사진에 사용."""
    n = len(images)
    if n == 0:
        return
    _photo_grid(slide, images, top, height, cols=n, rows=1, show_caption=show_caption)


def _points_panel(slide, points: List[dict], left, top, width, height, heading="핵심 보수 방법"):
    """"제목 1줄 + 설명 최대 2줄" 형태의 핵심 포인트 목록(최대 4개)을, 항목 수에
    맞춰 세로 공간을 균등 분배해 그린다 - 텍스트가 겹치거나 잘리지 않도록 항상
    포인트 개수 기준으로 행 높이를 계산한다(글자를 억지로 줄이지 않음)."""
    if not points:
        return
    pad = Mm(6)
    y = top + pad
    if heading:
        add_text(slide, left + pad, y, width - pad * 2, Mm(9), heading,
                  size=FONT_SIZE_SECTION_HEAD, color=COLOR_NAVY, bold=True)
        y += Mm(13)
    n = len(points)
    avail_h = height - (y - top) - pad
    row_h = avail_h / n
    for i, p in enumerate(points):
        ry = y + Emu(int(row_h * i))
        add_rect(slide, left + pad, ry + Mm(1.5), Mm(3), Mm(3), COLOR_GOLD)
        add_text(slide, left + pad + Mm(7), ry, width - pad * 2 - Mm(7), Mm(8),
                  f"{i+1}. {p.get('title','')}", size=FONT_SIZE_BODY, color=COLOR_NAVY, bold=True)
        add_text(slide, left + pad + Mm(7), ry + Mm(9), width - pad * 2 - Mm(7),
                  Emu(int(row_h - Mm(11))), p.get("desc", ""),
                  size=FONT_SIZE_CAPTION, color=COLOR_TEXT_BODY)


def _bullets_box(slide, bullets: List[str], left, top, width, height, bg=True):
    if not bullets:
        return
    if bg:
        add_rect(slide, left, top, width, height, COLOR_GRAY_LIGHT)
    pad = Mm(4)
    y = top + pad
    row_h = min(Mm(16), Emu(int((height - pad * 2) / max(len(bullets), 1))))
    for b in bullets:
        add_rect(slide, left + pad, y + Mm(3), Mm(3), Mm(3), COLOR_GOLD)
        add_text(slide, left + pad + Mm(7), y, width - pad * 2 - Mm(7), row_h, b,
                  size=FONT_SIZE_BODY, color=COLOR_TEXT_BODY)
        y += row_h


def _card_grid(slide, cards: List[dict], top, area_h, cols=2, accent=COLOR_NAVY):
    n = len(cards)
    if n == 0:
        return
    cols = min(cols, n)
    rows = (n + cols - 1) // cols
    gap = Mm(6)
    cell_w = (CONTENT_W - gap * (cols - 1)) / cols
    cell_h = (area_h - gap * (rows - 1)) / rows
    pad = Mm(4)
    label_h = Mm(9)
    cap_h = Mm(16)
    img_h = cell_h - label_h - cap_h - pad * 2
    for idx, card in enumerate(cards):
        r, c = divmod(idx, cols)
        x = MARGIN + Emu(int(cell_w + gap) * c)
        y = top + Emu(int(cell_h + gap) * r)
        add_rect(slide, x, y, Emu(int(cell_w)), Emu(int(cell_h)), COLOR_GRAY_LIGHT)
        add_rect(slide, x, y, Emu(int(cell_w)), Mm(2), accent)
        img = card.get("image")
        if img:
            add_picture_contain(slide, img.path, x + pad, y + pad, int(cell_w - pad * 2), int(img_h), frame=False)
        ty = y + pad + img_h
        add_text(slide, x + pad, ty, int(cell_w - pad * 2), label_h, card.get("label", ""),
                  size=FONT_SIZE_SECTION_HEAD, color=accent, bold=True)
        add_text(slide, x + pad, ty + label_h, int(cell_w - pad * 2), cap_h, card.get("caption", ""),
                  size=FONT_SIZE_CAPTION, color=COLOR_TEXT_SUB)


def _mark_pair_images(pair, images_by_id: Dict[str, object], page_no: int):
    for iid in [pair.before_image_id, pair.after_image_id] + list(pair.process_image_ids or []):
        im = images_by_id.get(iid)
        if im:
            im.selected_slide = page_no


# ------------------------------------------------------------------
# semantic_type별 전용 레이아웃
# ------------------------------------------------------------------

def render_hero(prs, page, page_no):
    """표지: 대표 사진 1장 크게 + 최소한의 문구, 여백 넓게."""
    slide = blank_slide(prs)
    img = page["images"][0] if page.get("images") else None
    if img:
        add_picture_contain(slide, img.path, 0, 0, SLIDE_WIDTH, Mm(190), frame=False)
        add_rect(slide, 0, Mm(178), SLIDE_WIDTH, Mm(2), COLOR_GOLD)
    add_rect(slide, 0, Mm(190), SLIDE_WIDTH, SLIDE_HEIGHT - Mm(190), COLOR_NAVY)
    add_text(slide, MARGIN, Mm(200), CONTENT_W, Mm(28), page["title"],
              size=Pt(30), color=COLOR_WHITE, bold=True)
    if page.get("subtitle"):
        add_text(slide, MARGIN, Mm(232), CONTENT_W, Mm(10), page["subtitle"],
                  size=Pt(14), color=(230, 232, 236))
    return slide


def render_reason_hero(prs, page, page_no):
    """공사가 필요한 이유: 대표 하자사진 1장 크게 + 보조 사진 + 핵심 문구 2~3개."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    top = Mm(27)
    lead = page.get("lead_image")
    support = page.get("support_images") or []
    lead_h = Mm(128) if support else Mm(150)
    if lead:
        add_picture_contain(slide, lead.path, MARGIN, top, int(CONTENT_W), int(lead_h))
    y = top + lead_h + Mm(6)
    if support:
        sup_h = Mm(36)
        _photo_row(slide, support, y, sup_h)
        y += sup_h + Mm(6)
    remaining_h = SLIDE_HEIGHT - y - Mm(12)
    _bullets_box(slide, page.get("bullets", []), MARGIN, y, CONTENT_W, remaining_h)
    return slide


def render_four_cards(prs, page, page_no):
    """주요 하자: 하자 유형별 사진+캡션 카드(2x2)."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    top = Mm(27)
    area_h = SLIDE_HEIGHT - top - Mm(12)
    _card_grid(slide, page.get("cards", []), top, area_h, cols=2, accent=COLOR_NAVY)
    return slide


def render_feature_cards(prs, page, page_no):
    """공법 특징: 특징별 사진+제목+설명 카드."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    top = Mm(27)
    area_h = SLIDE_HEIGHT - top - Mm(12)
    cards = page.get("cards", [])
    cols = 2 if len(cards) != 3 else 3
    _card_grid(slide, cards, top, area_h, cols=cols, accent=COLOR_GOLD)
    return slide


def render_image_text_split(prs, page, page_no):
    """보수 방법/공법 핵심: 좌측 55% 대표 시공사진 크게 + 우측 45% 핵심 포인트
    2~4개(제목 1줄+설명 최대 2줄) + 하단 보조사진 최대 2장(크게)."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    top = Mm(27)
    left_w = Emu(int(CONTENT_W * 0.55))
    right_w = CONTENT_W - left_w - Mm(6)
    lead = page.get("lead_image")
    support = (page.get("support_images") or [])[:2]
    support_h = Mm(55) if support else Mm(0)
    gap_before_support = Mm(6) if support else Mm(0)
    main_h = SLIDE_HEIGHT - top - Mm(12) - support_h - gap_before_support
    if lead:
        add_picture_contain(slide, lead.path, MARGIN, top, int(left_w), int(main_h))
    rx = MARGIN + left_w + Mm(6)
    add_rect(slide, rx, top, Emu(int(right_w)), int(main_h), COLOR_GRAY_LIGHT)
    points = page.get("points")
    if points:
        _points_panel(slide, points, rx, top, int(right_w), int(main_h))
    else:
        _bullets_box(slide, page.get("bullets", []), rx, top, int(right_w), int(main_h), bg=False)
    if support:
        _photo_row(slide, support, top + main_h + gap_before_support, support_h)
    return slide


def render_process_timeline(prs, page, page_no):
    """시공 순서: 단계별 사진 + 단계명 + 실제 원본 설명(공정 흐름 디자인)."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    top = Mm(27)
    area_h = SLIDE_HEIGHT - top - Mm(12)
    steps = page.get("steps", [])
    n = len(steps)
    if n == 0:
        return slide
    row_h = area_h / n
    img_w = Mm(58)
    badge_d = Mm(11)
    if n > 1:
        line_x = MARGIN + Emu(int(badge_d / 2)) - Mm(0.4)
        add_rect(slide, line_x, top + Mm(6), Mm(0.8), Emu(int(row_h * n - Mm(12))), COLOR_GRAY_LIGHT)
    for i, st in enumerate(steps):
        y = top + Emu(int(row_h * i))
        add_rect(slide, MARGIN, y + Mm(2), badge_d, badge_d, COLOR_NAVY)
        add_text(slide, MARGIN, y + Mm(2), badge_d, badge_d, str(st["step_no"]),
                  size=Pt(13), color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        img_x = MARGIN + badge_d + Mm(5)
        img_h = row_h - Mm(6)
        add_picture_contain(slide, st["image"].path, img_x, y + Mm(1), int(img_w), int(img_h), frame=False)
        text_x = img_x + img_w + Mm(6)
        text_w = CONTENT_W - (text_x - MARGIN)
        add_text(slide, text_x, y + Mm(2), int(text_w), Mm(9), st["label"],
                  size=FONT_SIZE_SECTION_HEAD, color=COLOR_NAVY, bold=True)
        add_text(slide, text_x, y + Mm(12), int(text_w), Emu(int(row_h - Mm(14))), st["desc"],
                  size=FONT_SIZE_BODY, color=COLOR_TEXT_BODY)
    return slide


def render_before_after(prs, page, page_no):
    """시공 전후 사례: 화면 대부분을 BEFORE/AFTER 큰 사진에 사용, 하단 짧은 설명."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    pair = page["pair"]
    images_by_id: Dict[str, object] = page["_images_by_id"]
    before = images_by_id.get(pair.before_image_id)
    after = images_by_id.get(pair.after_image_id)
    half_w = (CONTENT_W - Mm(4)) / 2
    top = Mm(29)
    ph = Mm(215)
    if before:
        add_picture_contain(slide, before.path, MARGIN, top, int(half_w), ph)
        add_rect(slide, MARGIN, top + ph + Mm(2), Emu(int(half_w)), Mm(9), COLOR_NAVY)
        add_text(slide, MARGIN, top + ph + Mm(2), Emu(int(half_w)), Mm(9), "시공 전",
                  size=Pt(13), color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if after:
        ax = MARGIN + Emu(int(half_w + Mm(4)))
        add_picture_contain(slide, after.path, ax, top, int(half_w), ph)
        add_rect(slide, ax, top + ph + Mm(2), Emu(int(half_w)), Mm(9), COLOR_GOLD)
        add_text(slide, ax, top + ph + Mm(2), Emu(int(half_w)), Mm(9), "시공 후",
                  size=Pt(13), color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    note_top = top + ph + Mm(15)
    bullets = page.get("bullets", [])
    if bullets:
        add_text(slide, MARGIN, note_top, CONTENT_W, Mm(14), bullets[0],
                  size=FONT_SIZE_BODY, color=COLOR_TEXT_BODY, align=PP_ALIGN.CENTER)
    return slide


def render_two_case_compare(prs, page, page_no):
    """여러 시공 전후 사례를 한 페이지에 압축 비교(최대 2사례)."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    images_by_id: Dict[str, object] = page["_images_by_id"]
    cases = page.get("cases", [])
    n = len(cases)
    if n == 0:
        return slide
    top = Mm(27)
    area_h = SLIDE_HEIGHT - top - Mm(12)
    row_h = area_h / n
    for i, c in enumerate(cases):
        y = top + Emu(int(row_h * i))
        pair = c["pair"]
        before = images_by_id.get(pair.before_image_id)
        after = images_by_id.get(pair.after_image_id)
        half_w = (CONTENT_W - Mm(4)) / 2
        img_h = row_h - Mm(20)
        add_text(slide, MARGIN, y, CONTENT_W, Mm(8), f"사례 {i+1}",
                  size=FONT_SIZE_SECTION_HEAD, color=COLOR_GOLD, bold=True)
        if before:
            add_picture_contain(slide, before.path, MARGIN, y + Mm(9), int(half_w), int(img_h), frame=False)
        if after:
            add_picture_contain(slide, after.path, MARGIN + Emu(int(half_w + Mm(4))), y + Mm(9),
                                  int(half_w), int(img_h), frame=False)
        add_text(slide, MARGIN, y + Mm(9) + img_h + Mm(1), CONTENT_W, Mm(9), c.get("note", ""),
                  size=FONT_SIZE_CAPTION, color=COLOR_TEXT_SUB)
    return slide


def render_effects_hero(prs, page, page_no):
    """기대 효과: 완료사진 1장 크게 + 핵심 효과 문구를 짧은 카드(칩)로."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    top = Mm(27)
    img = page.get("image") or (page["images"][0] if page.get("images") else None)
    img_h = Mm(145)
    if img:
        add_picture_contain(slide, img.path, MARGIN, top, int(CONTENT_W), int(img_h))
    y = top + img_h + Mm(8)
    bullets = page.get("bullets", [])
    chip_h = Mm(17)
    gap = Mm(5)
    for i, b in enumerate(bullets):
        by = y + Emu(int((chip_h + gap) * i))
        add_rect(slide, MARGIN, by, CONTENT_W, chip_h, COLOR_GRAY_LIGHT)
        add_rect(slide, MARGIN, by, Mm(3), chip_h, COLOR_GOLD)
        add_text(slide, MARGIN + Mm(8), by, CONTENT_W - Mm(12), chip_h, b,
                  size=FONT_SIZE_BODY, color=COLOR_TEXT_BODY, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def render_closing(prs, page, page_no):
    """마무리: 여백 중심, 문구 최소화, 회사정보 없음."""
    slide = blank_slide(prs)
    img = page["images"][0] if page.get("images") else None
    if img:
        add_picture_contain(slide, img.path, 0, 0, SLIDE_WIDTH, Mm(170), frame=False)
        add_rect(slide, 0, Mm(170), SLIDE_WIDTH, SLIDE_HEIGHT - Mm(170), COLOR_NAVY)
        text_top = Mm(170) + (SLIDE_HEIGHT - Mm(170)) / 2 - Mm(14)
    else:
        add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, COLOR_NAVY)
        text_top = SLIDE_HEIGHT / 2 - Mm(14)
    add_rect(slide, MARGIN, int(text_top) - Mm(8), CONTENT_W, Pt(2), COLOR_GOLD)
    add_text(slide, MARGIN, int(text_top), CONTENT_W, Mm(24), page["title"],
              size=Pt(19), color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    for i, b in enumerate(page.get("bullets", [])[:2]):
        add_text(slide, MARGIN, int(text_top) + Mm(24) + Emu(int(Mm(9) * i)), CONTENT_W, Mm(9), b,
                  size=Pt(12), color=(220, 222, 228), align=PP_ALIGN.CENTER)
    return slide


def _guidance_box(slide, text, left, top, width, height):
    """현장사진 페이지 하단의 공종별 공통 안내문구(문단형, 불릿 없음)."""
    if not text:
        return
    add_rect(slide, left, top, width, height, COLOR_GRAY_LIGHT)
    add_rect(slide, left, top, Mm(3), height, COLOR_GOLD)
    add_text(slide, left + Mm(8), top, width - Mm(14), height, text,
              size=FONT_SIZE_CAPTION, color=COLOR_TEXT_BODY, anchor=MSO_ANCHOR.MIDDLE)


def render_site_photo_gallery(prs, page, page_no):
    """현장사진: 사용자가 추가한 이 아파트의 실제 현장사진. 장수에 따라 배치를
    바꾸고, 억지로 작게 축소해 한 페이지에 전부 넣지 않는다(크게 보여주는 것이 원칙).
    개별 사진에는 어떤 캡션도 붙이지 않고(AI가 임의로 하자를 판단하지 않음),
    대신 공종 기준 공통 안내문구 하나를 하단에 표시한다."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    top = Mm(27)
    guidance_text = (page.get("bullets") or [None])[0]
    guidance_h = Mm(24) if guidance_text else Mm(0)
    gap_before_guidance = Mm(6) if guidance_text else Mm(0)
    area_h = SLIDE_HEIGHT - top - Mm(12) - guidance_h - gap_before_guidance
    images = page.get("images", [])
    n = len(images)
    if n == 0:
        return slide
    if n == 1:
        add_picture_contain(slide, images[0].path, MARGIN, top, int(CONTENT_W), int(area_h))
    elif n == 2:
        half_h = (area_h - Mm(4)) / 2
        add_picture_contain(slide, images[0].path, MARGIN, top, int(CONTENT_W), int(half_h))
        add_picture_contain(slide, images[1].path, MARGIN, top + half_h + Mm(4), int(CONTENT_W), int(half_h))
    elif n == 3:
        lead_h = area_h * 0.62
        add_picture_contain(slide, images[0].path, MARGIN, top, int(CONTENT_W), int(lead_h))
        sup_h = area_h - lead_h - Mm(4)
        _photo_row(slide, images[1:3], top + lead_h + Mm(4), sup_h, show_caption=False)
    elif n == 4:
        _photo_grid(slide, images, top, area_h, cols=2, rows=2, show_caption=False)
    else:
        lead_h = area_h * 0.55
        add_picture_contain(slide, images[0].path, MARGIN, top, int(CONTENT_W), int(lead_h))
        rest = images[1:6]
        sup_h = area_h - lead_h - Mm(4)
        cols = min(len(rest), 3) or 1
        rows = (len(rest) + cols - 1) // cols
        _photo_grid(slide, rest, top + lead_h + Mm(4), sup_h, cols=cols, rows=rows, show_caption=False)

    if guidance_text:
        _guidance_box(slide, guidance_text, MARGIN, top + area_h + gap_before_guidance,
                       CONTENT_W, guidance_h)
    return slide


def render_material_cards(prs, page, page_no):
    """사용 재료: 자재별 사진+명칭+원본 설명 카드."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    top = Mm(27)
    area_h = SLIDE_HEIGHT - top - Mm(12)
    cards = page.get("cards", [])
    cols = 2 if len(cards) != 3 else 3
    _card_grid(slide, cards, top, area_h, cols=cols, accent=COLOR_NAVY_LIGHT)
    return slide


def render_content_page(prs, page, page_no):
    """폴백 레이아웃(semantic_type이 없는 예외적인 경우에만 사용)."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    images = page.get("images", [])
    bullets = page.get("bullets", [])
    top = Mm(27)
    bottom_margin = Mm(12)
    total_h = SLIDE_HEIGHT - top - bottom_margin
    if images and bullets:
        text_h = min(Emu(int(total_h * 0.42)), Mm(9) * len(bullets) + Mm(8))
        text_h = max(text_h, Mm(20))
        img_h = total_h - text_h - Mm(4)
        _photo_grid(slide, images, top, img_h)
        _bullets_box(slide, bullets, MARGIN, top + img_h + Mm(4), CONTENT_W, text_h)
    elif images:
        _photo_grid(slide, images, top, total_h)
    else:
        _bullets_box(slide, bullets, MARGIN, top, CONTENT_W, total_h)
    return slide


# ------------------------------------------------------------------
# 추가 레이아웃 변형(디자인 시스템) - 같은 semantic_type이라도 페이지마다
# 다른 시각적 스타일을 골라 쓸 수 있도록 각 계열에 2~3개씩 변형을 마련한다.
# ------------------------------------------------------------------

GHOST_NUMBER_COLOR = (225, 228, 234)  # 큰 숫자 타이포를 위한 옅은 배경색 계열
_STAT_RE = re.compile(r"(\d[\d,]*\.?\d*)\s*(년|개|장|건|명|세대|층|%|m²|㎡|일|주|배)?")


def _find_stat_bullet(bullets: List[str]):
    """문구 중 실제로 숫자가 포함된 것이 있으면 (전체문구, 숫자, 단위)를 반환한다.
    원본에 없는 숫자를 새로 만들지 않고, 있는 경우에만 '숫자 강조' 레이아웃을 쓴다."""
    for b in bullets or []:
        m = _STAT_RE.search(b)
        if m and m.group(1):
            return b, m.group(1), (m.group(2) or "")
    return None


def _card_grid_minimal(slide, cards: List[dict], top, area_h, cols=2, accent=COLOR_NAVY):
    """카드형이지만 회색 박스 없이 사진 + 얇은 강조선 + 여백 위주(프리미엄/미니멀 톤)."""
    n = len(cards)
    if n == 0:
        return
    cols = min(cols, n)
    rows = (n + cols - 1) // cols
    gap = Mm(10)
    cell_w = (CONTENT_W - gap * (cols - 1)) / cols
    cell_h = (area_h - gap * (rows - 1)) / rows
    label_h = Mm(9)
    cap_h = Mm(16)
    img_h = cell_h - label_h - cap_h - Mm(5)
    for idx, card in enumerate(cards):
        r, c = divmod(idx, cols)
        x = MARGIN + Emu(int(cell_w + gap) * c)
        y = top + Emu(int(cell_h + gap) * r)
        img = card.get("image")
        if img:
            add_picture_contain(slide, img.path, x, y, int(cell_w), int(img_h), frame=False)
        ty = y + img_h + Mm(3)
        add_rect(slide, x, ty, Mm(10), Pt(2.2), accent)
        add_text(slide, x, ty + Mm(3), int(cell_w), label_h, card.get("label", ""),
                  size=FONT_SIZE_SECTION_HEAD, color=COLOR_TEXT_BODY, bold=True)
        add_text(slide, x, ty + Mm(3) + label_h, int(cell_w), cap_h, card.get("caption", ""),
                  size=FONT_SIZE_CAPTION, color=COLOR_TEXT_SUB)


def render_cards_grid_minimal(prs, page, page_no, accent=COLOR_NAVY):
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    top = Mm(27)
    area_h = SLIDE_HEIGHT - top - Mm(12)
    _card_grid_minimal(slide, page.get("cards", []), top, area_h, cols=2, accent=accent)
    return slide


def render_cards_list_numbered(prs, page, page_no, accent=COLOR_NAVY):
    """카드 목록을 그리드 대신 큰 옅은 숫자(ghost number) + 사진 + 텍스트의 세로
    리스트로 배치한다(Canva/Gamma류에서 흔한 번호 강조 타이포 기법)."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    top = Mm(27)
    area_h = SLIDE_HEIGHT - top - Mm(12)
    cards = page.get("cards", [])
    n = len(cards)
    if n == 0:
        return slide
    row_h = area_h / n
    num_w = Mm(22)
    img_w = Mm(50)
    for i, card in enumerate(cards):
        y = top + Emu(int(row_h * i))
        add_text(slide, MARGIN, y, num_w, row_h, f"{i+1:02d}", size=Pt(38),
                  color=GHOST_NUMBER_COLOR, bold=True)
        img = card.get("image")
        img_x = MARGIN + num_w
        if img:
            add_picture_contain(slide, img.path, img_x, y + Mm(2), int(img_w), int(row_h - Mm(4)), frame=False)
        text_x = img_x + img_w + Mm(6)
        text_w = CONTENT_W - (text_x - MARGIN)
        add_text(slide, text_x, y + Mm(4), int(text_w), Mm(9), card.get("label", ""),
                  size=FONT_SIZE_SECTION_HEAD, color=accent, bold=True)
        add_text(slide, text_x, y + Mm(14), int(text_w), Emu(int(row_h - Mm(16))), card.get("caption", ""),
                  size=FONT_SIZE_BODY, color=COLOR_TEXT_BODY)
        if i < n - 1:
            add_rect(slide, MARGIN, y + row_h - Mm(1), CONTENT_W, Pt(0.75), COLOR_GRAY_LIGHT)
    return slide


def render_material_cards_minimal(prs, page, page_no):
    return render_cards_list_numbered(prs, page, page_no, accent=COLOR_NAVY_LIGHT)


def render_split_right_image(prs, page, page_no):
    """image_text_split의 좌우 반전판: 핵심 포인트가 왼쪽 45%, 대표 사진이 오른쪽 55%."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    top = Mm(27)
    right_w = Emu(int(CONTENT_W * 0.55))
    left_w = CONTENT_W - right_w - Mm(6)
    lead = page.get("lead_image")
    support = (page.get("support_images") or [])[:2]
    support_h = Mm(55) if support else Mm(0)
    gap_before_support = Mm(6) if support else Mm(0)
    main_h = SLIDE_HEIGHT - top - Mm(12) - support_h - gap_before_support
    add_rect(slide, MARGIN, top, Emu(int(left_w)), int(main_h), COLOR_GRAY_LIGHT)
    points = page.get("points")
    if points:
        _points_panel(slide, points, MARGIN, top, int(left_w), int(main_h))
    else:
        _bullets_box(slide, page.get("bullets", []), MARGIN, top, int(left_w), int(main_h), bg=False)
    rx = MARGIN + left_w + Mm(6)
    if lead:
        add_picture_contain(slide, lead.path, rx, top, int(right_w), int(main_h))
    if support:
        _photo_row(slide, support, top + main_h + gap_before_support, support_h)
    return slide


def render_split_diagonal(prs, page, page_no):
    """대형 이미지 위에 핵심 포인트 카드가 겹쳐 올라간 에디토리얼 스타일 레이아웃."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    top = Mm(27)
    lead = page.get("lead_image")
    support = (page.get("support_images") or [])[:2]
    img_h = Mm(133)
    if lead:
        add_picture_contain(slide, lead.path, MARGIN, top, int(CONTENT_W), int(img_h))
    card_top = top + img_h - Mm(18)
    card_w = Emu(int(CONTENT_W * 0.62))
    card_h = Mm(85)
    add_rect(slide, MARGIN, card_top, card_w, card_h, COLOR_WHITE)
    add_rect(slide, MARGIN, card_top, Mm(3), card_h, COLOR_GOLD)
    points = page.get("points")
    if points:
        _points_panel(slide, points, MARGIN + Mm(3), card_top, int(CONTENT_W * 0.62) - Mm(3), int(card_h),
                       heading=None)
    else:
        _bullets_box(slide, page.get("bullets", []), MARGIN + Mm(6), card_top + Mm(4),
                      int(CONTENT_W * 0.62) - Mm(10), int(card_h) - Mm(8), bg=False)
    y = card_top + card_h + Mm(8)
    if support:
        sup_h = SLIDE_HEIGHT - y - Mm(12)
        _photo_row(slide, support, y, sup_h)
    return slide


def render_reason_diagonal(prs, page, page_no):
    """공사가 필요한 이유의 변형: 큰 사진 위에 보조 사진 + 문구 카드를 겹쳐 배치."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    top = Mm(27)
    lead = page.get("lead_image")
    support = page.get("support_images") or []
    img_h = Mm(155)
    if lead:
        add_picture_contain(slide, lead.path, MARGIN, top, int(CONTENT_W), int(img_h))
    card_top = top + img_h - Mm(18)
    card_h = Mm(95)
    add_rect(slide, MARGIN, card_top, CONTENT_W, card_h, COLOR_WHITE)
    add_rect(slide, MARGIN, card_top, Mm(3), card_h, COLOR_GOLD)
    if support:
        sup_h = Mm(34)
        _photo_row(slide, support, card_top + Mm(4), sup_h)
        _bullets_box(slide, page.get("bullets", []), MARGIN + Mm(6), card_top + sup_h + Mm(8),
                      int(CONTENT_W) - Mm(12), int(card_h - sup_h) - Mm(12), bg=False)
    else:
        _bullets_box(slide, page.get("bullets", []), MARGIN + Mm(6), card_top + Mm(4),
                      int(CONTENT_W) - Mm(12), int(card_h) - Mm(8), bg=False)
    return slide


def render_feature_magazine(prs, page, page_no):
    """공법 특징의 매거진 스타일: 왼쪽 대형 대표 사진 + 오른쪽 번호가 매겨진 특징 목록."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    top = Mm(27)
    cards = page.get("cards", [])
    if not cards:
        return slide
    lead_img = cards[0].get("image")
    left_w = Mm(78)
    main_h = SLIDE_HEIGHT - top - Mm(12)
    if lead_img:
        add_picture_contain(slide, lead_img.path, MARGIN, top, int(left_w), int(main_h))
    rx = MARGIN + left_w + Mm(8)
    right_w = CONTENT_W - left_w - Mm(8)
    row_h = main_h / len(cards)
    for i, card in enumerate(cards):
        y = top + Emu(int(row_h * i))
        add_text(slide, rx, y, Mm(16), row_h, f"{i+1:02d}", size=Pt(26), color=GHOST_NUMBER_COLOR, bold=True)
        tx = rx + Mm(18)
        tw = right_w - Mm(18)
        add_text(slide, tx, y + Mm(2), int(tw), Mm(9), card.get("label", ""),
                  size=FONT_SIZE_SECTION_HEAD, color=COLOR_GOLD, bold=True)
        add_text(slide, tx, y + Mm(12), int(tw), Emu(int(row_h - Mm(14))), card.get("caption", ""),
                  size=FONT_SIZE_BODY, color=COLOR_TEXT_BODY)
    return slide


def render_timeline_chevron(prs, page, page_no):
    """시공 순서의 가로형 변형: 단계를 2열 카드로 배치하고 화살표로 흐름을 표시한다."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    top = Mm(27)
    area_h = SLIDE_HEIGHT - top - Mm(12)
    steps = page.get("steps", [])
    n = len(steps)
    if n == 0:
        return slide
    cols = min(n, 2)
    rows = (n + cols - 1) // cols
    gap = Mm(10)
    cell_w = (CONTENT_W - gap * (cols - 1)) / cols
    cell_h = (area_h - gap * (rows - 1)) / rows
    img_h = cell_h * 0.5
    for i, st in enumerate(steps):
        r, c = divmod(i, cols)
        x = MARGIN + Emu(int(cell_w + gap) * c)
        y = top + Emu(int(cell_h + gap) * r)
        add_picture_contain(slide, st["image"].path, x, y, int(cell_w), int(img_h), frame=False)
        add_rect(slide, x, y + img_h + Mm(2), Mm(10), Pt(2.2), COLOR_GOLD)
        add_text(slide, x, y + img_h + Mm(5), int(cell_w), Mm(9), f"{st['step_no']:02d}  {st['label']}",
                  size=FONT_SIZE_SECTION_HEAD, color=COLOR_NAVY, bold=True)
        add_text(slide, x, y + img_h + Mm(15), int(cell_w), Emu(int(cell_h - img_h - Mm(17))), st["desc"],
                  size=FONT_SIZE_CAPTION, color=COLOR_TEXT_BODY)
        if c < cols - 1 and i + 1 < n:
            add_text(slide, x + cell_w, y + img_h / 2 - Mm(6), gap, Mm(12), "→",
                      size=Pt(18), color=COLOR_GOLD, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def render_before_after_ribbon(prs, page, page_no):
    """시공 전후 사례의 변형: 하단 바 대신 모서리 리본형 BEFORE/AFTER 라벨."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    pair = page["pair"]
    images_by_id: Dict[str, object] = page["_images_by_id"]
    before = images_by_id.get(pair.before_image_id)
    after = images_by_id.get(pair.after_image_id)
    half_w = (CONTENT_W - Mm(4)) / 2
    top = Mm(29)
    ph = Mm(205)
    if before:
        add_picture_contain(slide, before.path, MARGIN, top, int(half_w), ph, frame=False)
        add_rect(slide, MARGIN, top, Mm(30), Mm(9), COLOR_NAVY)
        add_text(slide, MARGIN, top, Mm(30), Mm(9), "BEFORE", size=Pt(11), color=COLOR_WHITE,
                  bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if after:
        ax = MARGIN + Emu(int(half_w + Mm(4)))
        add_picture_contain(slide, after.path, ax, top, int(half_w), ph, frame=False)
        add_rect(slide, ax, top, Mm(30), Mm(9), COLOR_GOLD)
        add_text(slide, ax, top, Mm(30), Mm(9), "AFTER", size=Pt(11), color=COLOR_WHITE,
                  bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    note_top = top + ph + Mm(8)
    bullets = page.get("bullets", [])
    if bullets:
        add_text(slide, MARGIN, note_top, CONTENT_W, Mm(14), bullets[0],
                  size=FONT_SIZE_BODY, color=COLOR_TEXT_BODY, align=PP_ALIGN.CENTER)
    return slide


def render_case_compare_grid(prs, page, page_no):
    """추가 사례 비교의 변형: 세로로 쌓지 않고 좌우 카드로 나란히 배치."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    images_by_id: Dict[str, object] = page["_images_by_id"]
    cases = page.get("cases", [])
    n = len(cases)
    if n == 0:
        return slide
    top = Mm(27)
    area_h = SLIDE_HEIGHT - top - Mm(12)
    cols = min(n, 2)
    rows = (n + cols - 1) // cols
    gap = Mm(8)
    cell_w = (CONTENT_W - gap * (cols - 1)) / cols
    cell_h = (area_h - gap * (rows - 1)) / rows
    for i, c in enumerate(cases):
        r, cc = divmod(i, cols)
        x = MARGIN + Emu(int(cell_w + gap) * cc)
        y = top + Emu(int(cell_h + gap) * r)
        pair = c["pair"]
        before = images_by_id.get(pair.before_image_id)
        after = images_by_id.get(pair.after_image_id)
        add_text(slide, x, y, int(cell_w), Mm(8), f"사례 {i+1}",
                  size=FONT_SIZE_SECTION_HEAD, color=COLOR_GOLD, bold=True)
        half_w = (cell_w - Mm(3)) / 2
        img_h = cell_h - Mm(18)
        if before:
            add_picture_contain(slide, before.path, x, y + Mm(9), int(half_w), int(img_h), frame=False)
        if after:
            add_picture_contain(slide, after.path, x + Emu(int(half_w + Mm(3))), y + Mm(9),
                                  int(half_w), int(img_h), frame=False)
        add_text(slide, x, y + Mm(9) + img_h + Mm(1), int(cell_w), Mm(8), c.get("note", ""),
                  size=FONT_SIZE_CAPTION, color=COLOR_TEXT_SUB)
    return slide


def render_effects_stat_highlight(prs, page, page_no):
    """기대 효과의 변형: 원본 문구에 실제 숫자가 있을 때만 그 숫자를 크게 강조한다
    (원본에 없는 숫자를 새로 만들지 않음 - _find_stat_bullet이 없으면 이 변형은 선택되지 않는다)."""
    slide = blank_slide(prs)
    header(slide, page["title"], page_no)
    top = Mm(27)
    bullets = page.get("bullets", [])
    stat = _find_stat_bullet(bullets)
    img = page.get("image") or (page["images"][0] if page.get("images") else None)
    left_w = Mm(75)
    main_h = SLIDE_HEIGHT - top - Mm(12)
    if img:
        add_picture_contain(slide, img.path, MARGIN, top, int(left_w), int(main_h))
    rx = MARGIN + left_w + Mm(8)
    rw = CONTENT_W - left_w - Mm(8)
    y = top
    remaining = bullets
    if stat:
        full_text, num, unit = stat
        add_text(slide, rx, y, int(rw), Mm(30), num + unit, size=Pt(46), color=COLOR_NAVY, bold=True)
        y += Mm(32)
        rest_label = full_text.replace(num, "", 1).strip(" :-·")
        if rest_label:
            add_text(slide, rx, y, int(rw), Mm(14), rest_label, size=FONT_SIZE_BODY, color=COLOR_TEXT_SUB)
            y += Mm(18)
        remaining = [b for b in bullets if b != full_text]
    chip_h = Mm(15)
    gap = Mm(4)
    for b in remaining:
        add_rect(slide, rx, y, int(rw), chip_h, COLOR_GRAY_LIGHT)
        add_rect(slide, rx, y, Mm(3), chip_h, COLOR_GOLD)
        add_text(slide, rx + Mm(7), y, int(rw) - Mm(10), chip_h, b,
                  size=FONT_SIZE_CAPTION, color=COLOR_TEXT_BODY, anchor=MSO_ANCHOR.MIDDLE)
        y += chip_h + gap
    return slide


def render_closing_quote(prs, page, page_no):
    """마무리의 변형: 상하 얇은 금색 룰 사이에 문구를 배치하는 절제된 인용구 스타일."""
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, COLOR_NAVY)
    center_y = SLIDE_HEIGHT / 2 - Mm(20)
    add_rect(slide, MARGIN + Mm(60), int(center_y) - Mm(10), CONTENT_W - Mm(120), Pt(1.5), COLOR_GOLD)
    add_text(slide, MARGIN, int(center_y), CONTENT_W, Mm(28), page["title"],
              size=Pt(22), color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    for i, b in enumerate(page.get("bullets", [])[:2]):
        add_text(slide, MARGIN, int(center_y) + Mm(30) + Emu(int(Mm(9) * i)), CONTENT_W, Mm(9), b,
                  size=Pt(12), color=(220, 222, 228), align=PP_ALIGN.CENTER)
    add_rect(slide, MARGIN + Mm(60), int(center_y) + Mm(50), CONTENT_W - Mm(120), Pt(1.5), COLOR_GOLD)
    return slide


SEMANTIC_RENDERERS = {
    "hero": render_hero,
    "site_photo_gallery": render_site_photo_gallery,
    "reason_hero": render_reason_hero,
    "four_cards": render_four_cards,
    "feature_cards": render_feature_cards,
    "material_cards": render_material_cards,
    "image_text_split": render_image_text_split,
    "process_timeline": render_process_timeline,
    "before_after": render_before_after,
    "two_case_compare": render_two_case_compare,
    "effects_hero": render_effects_hero,
    "closing": render_closing,
}

# semantic_type 하나당 2~3개의 시각적 변형(총 20개 이상)을 등록한다. 목록의 첫
# 항목이 기본값이며, 실제 선택은 generate_pptx_v2가 "직전 페이지와 다른 변형"을
# 우선하도록 고른다.
LAYOUT_CATALOG = {
    "hero": ["hero_classic"],
    "site_photo_gallery": ["site_gallery_standard"],
    "reason_hero": ["reason_hero_stacked", "reason_hero_diagonal"],
    "four_cards": ["cards_grid_classic", "cards_grid_minimal", "cards_list_numbered"],
    "image_text_split": ["split_left_image", "split_right_image", "split_diagonal"],
    "feature_cards": ["feature_cards_grid", "feature_cards_magazine"],
    "material_cards": ["material_cards_grid", "material_cards_minimal"],
    "process_timeline": ["timeline_vertical", "timeline_chevron"],
    "before_after": ["before_after_split", "before_after_ribbon"],
    "two_case_compare": ["case_compare_stacked", "case_compare_grid"],
    "effects_hero": ["effects_hero_standard", "effects_stat_highlight"],
    "closing": ["closing_minimal", "closing_quote"],
}

VARIANT_RENDERERS = {
    "hero_classic": render_hero,
    "site_gallery_standard": render_site_photo_gallery,
    "reason_hero_stacked": render_reason_hero,
    "reason_hero_diagonal": render_reason_diagonal,
    "cards_grid_classic": render_four_cards,
    "cards_grid_minimal": render_cards_grid_minimal,
    "cards_list_numbered": render_cards_list_numbered,
    "split_left_image": render_image_text_split,
    "split_right_image": render_split_right_image,
    "split_diagonal": render_split_diagonal,
    "feature_cards_grid": render_feature_cards,
    "feature_cards_magazine": render_feature_magazine,
    "material_cards_grid": render_material_cards,
    "material_cards_minimal": render_material_cards_minimal,
    "timeline_vertical": render_process_timeline,
    "timeline_chevron": render_timeline_chevron,
    "before_after_split": render_before_after,
    "before_after_ribbon": render_before_after_ribbon,
    "case_compare_stacked": render_two_case_compare,
    "case_compare_grid": render_case_compare_grid,
    "effects_hero_standard": render_effects_hero,
    "effects_stat_highlight": render_effects_stat_highlight,
    "closing_minimal": render_closing,
    "closing_quote": render_closing_quote,
}


def _eligible_variants(stype: str, page: dict) -> List[str]:
    """페이지 내용상 실제로 사용 가능한 변형만 남긴다(예: 원본 문구에 숫자가
    없으면 '숫자 강조' 변형은 후보에서 제외 - 없는 사실을 만들어내지 않기 위함)."""
    candidates = list(LAYOUT_CATALOG.get(stype, [stype]))
    if stype == "effects_hero" and not _find_stat_bullet(page.get("bullets", [])):
        candidates = [v for v in candidates if v != "effects_stat_highlight"]
    return candidates or list(LAYOUT_CATALOG.get(stype, [stype]))


def generate_pptx_v2(pages: List[dict], images_by_id: Dict[str, object], out_path: str,
                       log_fn=None) -> str:
    """PPTX를 생성해 out_path에 저장한다. 저장 직후 반드시 파일이 실제로 존재하고
    크기가 0보다 큰지 확인한 뒤 로그를 남긴다 - "저장했다고 believe"하는 대신
    os.path.exists()로 실제 결과를 확인한다. 저장이 실패했으면 예외를 던져
    호출자(pipeline.py)가 "정상 완료"로 처리하지 못하게 한다."""
    def _log(msg):
        if log_fn:
            log_fn(msg)
        else:
            print(msg, flush=True)

    prs = new_presentation()
    page_no = 0
    prev_variant = None  # 직전 페이지에 사용한 레이아웃 변형(같은 변형이 연속되지 않도록 추적)
    variant_log = []

    for page in pages:
        t = page["type"]
        stype = page.get("semantic_type")

        if t == "cover":
            render_hero(prs, page, None)
            for im in page.get("images", []):
                im.selected_slide = 1
            prev_variant = "hero_classic"
            continue

        page_no += 1

        if stype in ("before_after", "two_case_compare"):
            page["_images_by_id"] = images_by_id

        candidates = _eligible_variants(stype, page) if stype else []
        variant = next((v for v in candidates if v != prev_variant), candidates[0]) if candidates else None
        renderer = VARIANT_RENDERERS.get(variant) if variant else None

        if renderer:
            renderer(prs, page, page_no)
            prev_variant = variant
            variant_log.append(f"{page_no}:{page.get('title','')}={variant}")
        elif stype in SEMANTIC_RENDERERS:
            SEMANTIC_RENDERERS[stype](prs, page, page_no)
            prev_variant = stype
            variant_log.append(f"{page_no}:{page.get('title','')}={stype}(fallback)")
        else:
            render_content_page(prs, page, page_no)
            prev_variant = "content_fallback"
            variant_log.append(f"{page_no}:{page.get('title','')}=content_fallback")

        if stype == "before_after":
            _mark_pair_images(page["pair"], images_by_id, page_no)
        elif stype == "two_case_compare":
            for c in page.get("cases", []):
                _mark_pair_images(c["pair"], images_by_id, page_no)
        else:
            for im in page.get("images", []):
                im.selected_slide = page_no

    _log("[generator2] 페이지별 레이아웃 변형 선택 결과 - " + " / ".join(variant_log))
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    _log(f"[generator2] Presentation.save() 호출 직전 - 총 {len(prs.slides)}슬라이드, 저장 경로={out_path}")
    prs.save(out_path)

    saved_ok = os.path.exists(out_path)
    size = os.path.getsize(out_path) if saved_ok else 0
    if not saved_ok or size == 0:
        _log(f"[generator2] PPT 저장 실패 - 경로={out_path}, 존재={saved_ok}, 크기={size} bytes")
        raise RuntimeError(
            f"PPTX 저장에 실패했습니다(파일이 생성되지 않았거나 비어 있음): {out_path}"
        )
    _log(f"[generator2] PPT 저장 완료 - 경로={out_path}, 파일크기={size:,} bytes")
    return out_path
