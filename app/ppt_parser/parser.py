# -*- coding: utf-8 -*-
"""
기존 PPT 파일을 열어 슬라이드/도형 단위로 분석한다.
원본 파일은 절대 수정하지 않으며, 항상 임시 작업 폴더에 복사한 뒤 읽기 전용으로 연다.
"""
import os
import shutil
import uuid
from typing import List, Tuple

from pptx import Presentation
from pptx.util import Emu

from app.utils.models import SlideRecord, TextRun


def make_workcopy(src_path: str, temp_dir: str) -> str:
    """원본을 절대 건드리지 않기 위해 임시 폴더에 복사본을 만든다."""
    os.makedirs(temp_dir, exist_ok=True)
    ext = os.path.splitext(src_path)[1]
    dst = os.path.join(temp_dir, f"src_{uuid.uuid4().hex[:8]}{ext}")
    shutil.copy2(src_path, dst)
    return dst


def _iter_shapes(shapes):
    """그룹 도형 내부까지 재귀적으로 순회한다."""
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                yield from _iter_shapes(shape.shapes)
            except Exception:
                pass


def _extract_text_from_shape(shape) -> str:
    texts = []
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs)
                if line.strip():
                    texts.append(line.strip())
    except Exception:
        pass
    try:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text.strip())
    except Exception:
        pass
    return "\n".join(texts)


def parse_presentation(work_path: str, source_label: str) -> Tuple[List[SlideRecord], List[TextRun], "Presentation"]:
    """PPT 파일을 열어 슬라이드 레코드/텍스트 런 목록을 반환한다.
    이미지 자체 추출은 image_extractor 모듈이 담당한다 (Picture shape만 이 단계에서 식별).
    """
    prs = Presentation(work_path)
    slide_w, slide_h = prs.slide_width, prs.slide_height
    is_portrait = slide_h >= slide_w

    slide_records: List[SlideRecord] = []
    text_runs: List[TextRun] = []

    for idx, slide in enumerate(prs.slides):
        all_texts = []
        title_text = ""
        has_before = False
        has_after = False

        shapes = list(_iter_shapes(slide.shapes))

        # 제목 플레이스홀더 우선 탐색
        for shape in shapes:
            try:
                if shape.is_placeholder and shape.placeholder_format.type is not None:
                    ph_type = shape.placeholder_format.type
                    if ph_type in (13, 1) or "title" in str(ph_type).lower():  # TITLE류
                        t = _extract_text_from_shape(shape)
                        if t:
                            title_text = t
            except Exception:
                pass

        for shape in shapes:
            text = _extract_text_from_shape(shape)
            if text:
                all_texts.append(text)
                text_runs.append(TextRun(
                    source_file=source_label,
                    slide_index=idx,
                    shape_name=shape.name or "",
                    text=text,
                    is_title=(text == title_text),
                ))
                low = text.replace(" ", "").lower()
                if "시공전" in low or "before" in low or low.strip() in ("전",):
                    has_before = True
                if "시공후" in low or "after" in low or "완료" in low:
                    has_after = True

        joined = "\n".join(all_texts)
        if not title_text and all_texts:
            title_text = all_texts[0][:40]

        slide_records.append(SlideRecord(
            source_file=source_label,
            index=idx,
            is_portrait=is_portrait,
            title_text=title_text,
            all_text=joined,
            image_ids=[],  # image_extractor가 채운다
            has_before_after_labels=(has_before and has_after),
        ))

    return slide_records, text_runs, prs
