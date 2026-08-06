# -*- coding: utf-8 -*-
"""
입주민용 공종 설명자료 PPT 자동 제작 프로그램 - 파이프라인 오케스트레이터 + CLI.

사용법:
    python -m app.main --apt "행복아파트" --work 재도장 --output ./output file1.pptx file2.pptx [file3.pptx]
"""
import argparse
import datetime
import os
import subprocess
import sys
import tempfile
import traceback
from typing import Callable, List, Optional

from app.anonymizer.anonymizer import build_blacklist, evaluate_image
from app.content_selector import selector as cs
from app.content_selector.text_analyzer import (
    attach_real_captions, build_content_library, build_text_candidates,
)
from app.image_classifier.classifier import classify_all, dedupe_images, detect_before_after_pairs
from app.image_extractor.extractor import extract_images
from app.ppt_generator.generator import generate_pptx
from app.ppt_parser.parser import make_workcopy, parse_presentation
from app.slide_planner.planner import build_content_plan, plan_slides
from app.utils import debug_dump as dd
from app.validator.validator import validate_and_fix

STAGES = [
    "PPT 분석 중", "이미지 추출 중", "문구 분석 중", "기존 정보 제거 중",
    "사진 분류 중", "슬라이드 구성 중", "중간 산출물 검증 중",
    "PowerPoint 생성 중", "검수 중", "완료",
]


def _log(logs: List[str], msg: str):
    ts = datetime.datetime.now().strftime("%H:%M:%S")
    logs.append(f"[{ts}] {msg}")


def _report(cb: Optional[Callable[[str], None]], stage: str):
    if cb:
        cb(stage)


def convert_to_pdf(pptx_path: str, out_dir: str) -> Optional[str]:
    try:
        subprocess.run(
            ["soffice", "--headless", "--norestore", "--convert-to", "pdf", "--outdir", out_dir, pptx_path],
            check=True, timeout=180, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
        )
        pdf_path = os.path.join(out_dir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
        return pdf_path if os.path.exists(pdf_path) else None
    except Exception:
        return None


def build_preview_image(pdf_path: str, out_png: str, cols: int = 3) -> Optional[str]:
    try:
        import fitz  # PyMuPDF
        from PIL import Image

        doc = fitz.open(pdf_path)
        thumbs = []
        for page in doc:
            pix = page.get_pixmap(matrix=fitz.Matrix(0.6, 0.6))
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            thumbs.append(img)
        if not thumbs:
            return None
        rows = (len(thumbs) + cols - 1) // cols
        tw, th = thumbs[0].size
        pad = 10
        grid = Image.new("RGB", (cols * tw + (cols + 1) * pad, rows * th + (rows + 1) * pad), "white")
        for i, im in enumerate(thumbs):
            r, c = divmod(i, cols)
            x = pad + c * (tw + pad)
            y = pad + r * (th + pad)
            grid.paste(im, (x, y))
        grid.save(out_png)
        return out_png
    except Exception:
        return None


def run_pipeline(apartment_name: str, work_type: str, input_paths: List[str], output_dir: str,
                  progress_cb: Optional[Callable[[str], None]] = None) -> dict:
    logs: List[str] = []
    _log(logs, f"입력 검증 시작 - 아파트명='{apartment_name}', 공종='{work_type}', 파일수={len(input_paths)}")

    if not apartment_name or not apartment_name.strip():
        raise ValueError("새 아파트명을 입력해야 합니다.")
    if len(input_paths) < 2 or len(input_paths) > 3:
        raise ValueError("기존 PPT 파일은 2개 또는 3개를 입력해야 합니다.")
    for p in input_paths:
        if not os.path.exists(p):
            raise FileNotFoundError(f"파일을 찾을 수 없습니다: {p}")

    os.makedirs(output_dir, exist_ok=True)
    temp_root = tempfile.mkdtemp(prefix="resident_ppt_")
    _log(logs, f"임시 작업 폴더 생성: {temp_root}")

    # 1. PPT 분석
    _report(progress_cb, STAGES[0])
    all_slides, all_texts, source_files = [], [], []
    prs_objects = []
    for i, path in enumerate(input_paths):
        label = f"입력파일{i+1}"
        source_files.append(label)
        work_copy = make_workcopy(path, temp_root)
        slides, texts, prs = parse_presentation(work_copy, label)
        all_slides.extend(slides)
        all_texts.extend(texts)
        prs_objects.append((label, prs))
        _log(logs, f"{label}({os.path.basename(path)}) 분석 완료 - 슬라이드 {len(slides)}개, 텍스트 {len(texts)}건")

    total_slide_count = len(all_slides)
    _log(logs, f"[집계] 입력 PPT 전체 슬라이드 수: {total_slide_count}장")

    # 2. 이미지 추출
    _report(progress_cb, STAGES[1])
    all_images = []
    for label, prs in prs_objects:
        slides_for_file = [s for s in all_slides if s.source_file == label]
        imgs = extract_images(prs, label, slides_for_file, temp_root, run_ocr=True)
        all_images.extend(imgs)
        _log(logs, f"[집계] {label} 발견 이미지(Picture 도형) 수: {len(imgs)}장 (추출 성공 {len(imgs)}장)")
    _log(logs, f"[집계] 전체 추출 성공 이미지 수: {len(all_images)}장")

    # 3. 문구 추출 및 용도 분석 (실제 원본 문구 재사용을 위한 콘텐츠 라이브러리 구성)
    _report(progress_cb, STAGES[2])
    blacklist = build_blacklist(all_texts, all_images, apartment_name)
    text_candidates = build_text_candidates(all_texts, blacklist)
    usable_text_count = sum(1 for c in text_candidates if c.usable)
    _log(logs, f"[집계] 전체 텍스트 조각 {len(text_candidates)}개 중 재사용 가능 {usable_text_count}개")
    content_library = build_content_library(text_candidates)
    for purpose, phrases in content_library.items():
        _log(logs, f"  · {purpose}: {len(phrases)}건 확보")

    # 4. 기존 정보 제거 (블랙리스트 구성 + 배제 판정)
    _report(progress_cb, STAGES[3])
    _log(logs, f"블랙리스트 구성 완료 - 아파트명후보 {len(blacklist.apt_names)}, "
               f"회사명후보 {len(blacklist.companies)}, 전화 {len(blacklist.phones)}, "
               f"이메일 {len(blacklist.emails)}, URL {len(blacklist.urls)}")
    for img in all_images:
        evaluate_image(img, blacklist)
    banned_count = sum(1 for i in all_images if i.banned)
    banned_reason_tally = {}
    for i in all_images:
        for r in i.banned_reasons:
            banned_reason_tally[r] = banned_reason_tally.get(r, 0) + 1
    _log(logs, f"[집계] 민감정보(회사/현장 식별) 추정으로 배제된 이미지: {banned_count}장")
    for reason, cnt in banned_reason_tally.items():
        _log(logs, f"  · 배제 사유 '{reason}': {cnt}장")

    # 5. 사진 분류 + 중복 제거 + 전후 관계 탐지
    _report(progress_cb, STAGES[4])
    classify_all(all_images)
    dedupe_images(all_images)
    dup_count = sum(1 for i in all_images if i.is_duplicate_of)
    _log(logs, f"[집계] 중복/유사 사진으로 판단되어 제외된 이미지: {dup_count}장")
    detect_before_after_pairs(all_images, all_slides)
    attach_real_captions(all_images, text_candidates)

    images_by_id = {i.id: i for i in all_images}
    usable_images = [i for i in all_images if not i.banned and not i.is_duplicate_of]
    _log(logs, f"[집계] 최종 사용 가능(배제/중복 제외) 이미지: {len(usable_images)}장")
    cat_tally = {}
    for i in usable_images:
        cat_tally[i.category] = cat_tally.get(i.category, 0) + 1
    for cat, cnt in sorted(cat_tally.items(), key=lambda x: -x[1]):
        _log(logs, f"  · 분류 '{cat}': {cnt}장")

    # 6. 슬라이드 구성 (콘텐츠 선별 + 플랜) — 카테고리별로 뽑되, 다른 섹션에서 이미
    #    사용한 사진은 제외하여 서로 다른 사진을 최대한 많이 활용한다.
    _report(progress_cb, STAGES[5])
    used_ids = set()
    cover_id = cs.select_cover_image(all_images)
    if cover_id:
        used_ids.add(cover_id)

    defect_ids = cs.select_defect_images(all_images, exclude_ids=used_ids)
    used_ids.update(defect_ids)
    method_ids = cs.select_method_overview_images(all_images, exclude_ids=used_ids)
    used_ids.update(method_ids)
    feature_ids = cs.select_feature_images(all_images, exclude_ids=used_ids)
    used_ids.update(feature_ids)
    effect_ids = cs.select_effect_images(all_images, exclude_ids=used_ids)
    used_ids.update(effect_ids)

    process_steps = cs.build_process_steps(all_images, all_texts, source_files,
                                             text_candidates=text_candidates, exclude_ids=used_ids)
    for s in process_steps:
        used_ids.update(s.image_ids)

    ba_pairs = cs.build_before_after_cases(all_images)
    for p in ba_pairs:
        used_ids.add(p.before_image_id)
        used_ids.add(p.after_image_id)
        used_ids.update(p.process_image_ids)

    gallery_ids = cs.select_gallery_images(all_images, exclude_ids=used_ids, max_count=40)
    used_ids.update(gallery_ids)

    total_used = len(used_ids)
    _log(logs, f"[집계] 결과물 후보로 최종 선정된 서로 다른 이미지 수: {total_used}장 "
               f"(하자 {len(defect_ids)} / 공법개요 {len(method_ids)} / 특징 {len(feature_ids)} / "
               f"효과 {len(effect_ids)} / 시공순서 {sum(len(s.image_ids) for s in process_steps)} / "
               f"전후사례 {sum(2 + len(p.process_image_ids) for p in ba_pairs)} / 갤러리 {len(gallery_ids)})")

    method_items = cs.build_method_overview_items(content_library)
    feature_items = cs.build_feature_items(content_library)
    effect_items = cs.build_effect_items(content_library)
    reason_note = cs.build_reason_note(content_library)
    case_notes = cs.build_case_notes(content_library)

    warnings = []
    if not ba_pairs:
        warnings.append("전·후 관계가 명확한 사진 쌍을 찾지 못해 유사 시공 사례 페이지가 생략되었습니다.")
    if not process_steps:
        warnings.append("공법 순서를 확인할 수 있는 자료가 부족하여 시공 순서 페이지가 생략되었습니다.")
    needs_confirm = [s.name for s in process_steps if s.needs_user_confirmation]
    if needs_confirm:
        warnings.append(f"입력 자료 간 공정 순서 표기가 달라 사용자 확인이 필요한 단계: {', '.join(needs_confirm)}")
    if total_used < 30:
        warnings.append(f"입력 PPT에서 확보한 사용 가능 사진이 {total_used}장으로, 목표치(30장 이상)보다 적습니다. "
                          "원본 PPT의 사진 절대량 또는 회사/현장 정보 노출로 인한 배제가 원인일 수 있습니다.")

    for w in warnings:
        _log(logs, f"[경고] {w}")

    content_plan = build_content_plan(
        apartment_name, work_type, cover_id, defect_ids, method_ids, feature_ids,
        process_steps, ba_pairs, effect_ids, warnings,
        gallery_ids=gallery_ids, method_items=method_items, feature_items=feature_items,
        effect_items=effect_items, reason_note=reason_note, case_notes=case_notes,
    )
    slide_plan = plan_slides(content_plan)
    _log(logs, f"슬라이드 구성 완료 - 총 {len(slide_plan)}페이지 계획")

    # 7. 중간 산출물 저장 및 검증 (PPTX 생성 이전에 먼저 저장)
    _report(progress_cb, STAGES[6])
    safe_name = "".join(c for c in apartment_name if c not in '\\/:*?"<>|').strip() or "새아파트"
    debug_dir = os.path.join(output_dir, f"{safe_name}_중간산출물")
    os.makedirs(debug_dir, exist_ok=True)
    dd.dump_image_extraction_folder(all_images, debug_dir)
    dd.dump_image_classification_csv(all_images, os.path.join(debug_dir, "이미지분류.csv"))
    dd.dump_final_candidates_csv(images_by_id, {
        "표지": [cover_id] if cover_id else [],
        "하자": defect_ids, "공법개요": method_ids, "공법특징": feature_ids,
        "기대효과": effect_ids, "갤러리": gallery_ids,
        "시공순서": [iid for s in process_steps for iid in s.image_ids],
        "전후사례": [iid for p in ba_pairs for iid in
                    ([p.before_image_id, p.after_image_id] + list(p.process_image_ids or []))],
    }, os.path.join(debug_dir, "최종후보이미지목록.csv"))
    dd.dump_slide_plan_json(slide_plan, os.path.join(debug_dir, "슬라이드계획.json"))
    dd.dump_extracted_text_csv(text_candidates, os.path.join(debug_dir, "extracted_text.csv"))
    dd.dump_content_library_json(content_library, os.path.join(debug_dir, "content_library.json"))
    dd.dump_slide_content_mapping_csv(slide_plan, images_by_id,
                                        os.path.join(debug_dir, "slide_content_mapping.csv"))
    _log(logs, f"중간 산출물 저장 완료: {debug_dir}")

    # 8. PPTX 생성
    _report(progress_cb, STAGES[7])
    out_pptx = os.path.join(output_dir, f"{safe_name}_{work_type}_입주민설명자료.pptx")
    generate_pptx(slide_plan, images_by_id, out_pptx)

    # 실제로 최종 PPTX 안에 삽입된 이미지 수를 재확인(선정 목록이 아닌 실물 검증)
    from pptx import Presentation as _Presentation
    _verify_prs = _Presentation(out_pptx)
    inserted_count = 0

    def _count_pics(shapes):
        n = 0
        for sh in shapes:
            if sh.shape_type == 13:
                n += 1
            if sh.shape_type == 6:
                n += _count_pics(sh.shapes)
        return n
    for _slide in _verify_prs.slides:
        inserted_count += _count_pics(_slide.shapes)
    _log(logs, f"[집계] 실제 생성된 PPT에 삽입된 이미지 수(재오픈 검증): {inserted_count}장 / "
               f"전체 {len(_verify_prs.slides)}페이지")
    _log(logs, f"PowerPoint 파일 생성 완료: {out_pptx}")

    # 9. 검수
    _report(progress_cb, STAGES[8])
    report = validate_and_fix(out_pptx, blacklist, apartment_name)
    _log(logs, "자동 검수 완료")

    report_path = os.path.join(output_dir, f"{safe_name}_검수결과.txt")
    with open(report_path, "w", encoding="utf-8") as f:
        f.write(report.as_text())

    # PDF / 미리보기 생성 (LibreOffice가 없는 환경에서는 건너뜀)
    pdf_path = convert_to_pdf(out_pptx, output_dir)
    if pdf_path:
        wanted_pdf = os.path.join(output_dir, f"{safe_name}_{work_type}_입주민설명자료.pdf")
        if pdf_path != wanted_pdf:
            os.replace(pdf_path, wanted_pdf)
            pdf_path = wanted_pdf
        _log(logs, f"PDF 미리보기 생성 완료: {pdf_path}")
        preview_png = os.path.join(output_dir, f"{safe_name}_미리보기.png")
        if build_preview_image(pdf_path, preview_png):
            _log(logs, f"슬라이드 미리보기 이미지 생성 완료: {preview_png}")
        else:
            preview_png = None
            _log(logs, "미리보기 이미지 생성 실패(건너뜀)")
    else:
        preview_png = None
        _log(logs, "LibreOffice를 찾을 수 없어 PDF/미리보기 생성을 건너뛰었습니다.")

    log_path = os.path.join(output_dir, f"{safe_name}_처리로그.txt")
    with open(log_path, "w", encoding="utf-8") as f:
        f.write("\n".join(logs))

    _report(progress_cb, STAGES[9])

    return {
        "pptx": out_pptx,
        "pdf": pdf_path,
        "preview_png": preview_png,
        "log": log_path,
        "validation_report": report_path,
        "debug_dir": debug_dir,
        "inserted_image_count": inserted_count,
        "total_usable_image_count": total_used,
        "warnings": warnings,
        "validation": report,
        "temp_dir": temp_root,
    }


def main():
    parser = argparse.ArgumentParser(description="입주민 설명자료 자동 제작")
    parser.add_argument("files", nargs="+", help="기존 PPT 파일 2~3개")
    parser.add_argument("--apt", required=True, help="새 아파트명")
    parser.add_argument("--work", default="재도장", choices=["재도장", "방수", "보수·보강", "아스콘", "기타"])
    parser.add_argument("--output", default="./output", help="결과물 저장 폴더")
    args = parser.parse_args()

    def cb(stage):
        print(f"  >> {stage}")

    try:
        result = run_pipeline(args.apt, args.work, args.files, args.output, progress_cb=cb)
    except Exception as e:
        print(f"[오류] {e}")
        traceback.print_exc()
        sys.exit(1)

    print("\n=== 완료 ===")
    for k in ("pptx", "pdf", "preview_png", "log", "validation_report"):
        print(f"{k}: {result.get(k)}")
    if result["warnings"]:
        print("\n[확인 필요]")
        for w in result["warnings"]:
            print(f" - {w}")


if __name__ == "__main__":
    main()
