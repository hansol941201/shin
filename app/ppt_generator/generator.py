# -*- coding: utf-8 -*-
"""
계획된 슬라이드 구조를 실제 PowerPoint(.pptx) 파일로 렌더링한다.
- A4 세로형, 흰 배경 + 남색 주색상 + 회색/금색 포인트
- 사진은 원본 비율을 유지한 채 배치(찌그러뜨리지 않음)
"""
import os
from typing import Dict, List, Optional

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Mm, Pt

from app.utils.config import (
    SLIDE_WIDTH, SLIDE_HEIGHT, COLOR_WHITE, COLOR_NAVY, COLOR_NAVY_LIGHT,
    COLOR_GRAY_LIGHT, COLOR_GRAY_MID, COLOR_GOLD, COLOR_TEXT_BODY, COLOR_TEXT_SUB,
    FONT_PRIMARY, FONT_SIZE_COVER_TITLE, FONT_SIZE_COVER_SUB, FONT_SIZE_PAGE_TITLE,
    FONT_SIZE_SECTION_HEAD, FONT_SIZE_BODY, FONT_SIZE_CAPTION, FONT_SIZE_FOOTER,
    NEUTRAL_CAPTIONS,
)
from app.utils.models import ImageAsset

MARGIN = Mm(12)
CONTENT_W = SLIDE_WIDTH - MARGIN * 2


def _rgb(t):
    return RGBColor(*t)


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def _blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # 배경 흰색 고정
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(COLOR_WHITE)
    return slide


def _set_font(run, size=FONT_SIZE_BODY, color=COLOR_TEXT_BODY, bold=False, name=FONT_PRIMARY):
    run.font.size = size
    run.font.bold = bold
    run.font.name = name
    run.font.color.rgb = _rgb(color)
    try:
        rPr = run.font._rPr
        ea = rPr.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}ea',
                              {'typeface': name})
        rPr.append(ea)
    except Exception:
        pass


def add_rect(slide, left, top, width, height, color, line=False):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    if not line:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, size=FONT_SIZE_BODY, color=COLOR_TEXT_BODY,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, name=FONT_PRIMARY,
             line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        _set_font(run, size=size, color=color, bold=bold, name=name)
    return box


def _fit_box(img_w, img_h, max_w, max_h):
    """비율을 유지하며 max_w x max_h 안에 들어가는 크기를 계산한다."""
    ratio = min(max_w / img_w, max_h / img_h)
    return int(img_w * ratio), int(img_h * ratio)


def add_picture_contain(slide, path, left, top, max_w, max_h, frame=True):
    """가로세로 비율을 유지한 채 지정 영역 안에 이미지를 배치(중앙 정렬)."""
    try:
        with Image.open(path) as im:
            iw, ih = im.size
    except Exception:
        iw, ih = max_w, max_h
    w, h = _fit_box(iw, ih, max_w, max_h)
    x = left + (max_w - w) // 2
    y = top + (max_h - h) // 2
    if frame:
        pad = Emu(9525 * 2)
        add_rect(slide, x - pad, y - pad, w + pad * 2, h + pad * 2, COLOR_GRAY_LIGHT)
    pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
    return pic, (x, y, w, h)


def _header(slide, title_text, page_no=None):
    add_rect(slide, 0, 0, SLIDE_WIDTH, Mm(20), COLOR_NAVY)
    add_rect(slide, 0, Mm(20), SLIDE_WIDTH, Pt(2.2), COLOR_GOLD)
    add_text(slide, MARGIN, Mm(3.5), CONTENT_W, Mm(13), title_text,
              size=FONT_SIZE_PAGE_TITLE, color=COLOR_WHITE, bold=True,
              anchor=MSO_ANCHOR.MIDDLE)
    if page_no is not None:
        add_text(slide, SLIDE_WIDTH - Mm(30), SLIDE_HEIGHT - Mm(12), Mm(20), Mm(8),
                  f"{page_no:02d}", size=FONT_SIZE_FOOTER, color=COLOR_GRAY_MID,
                  align=PP_ALIGN.RIGHT)


def _neutral_caption(idx=0):
    return NEUTRAL_CAPTIONS[idx % len(NEUTRAL_CAPTIONS)]


# ------------------------------------------------------------------
# 슬라이드 유형별 렌더러
# ------------------------------------------------------------------

def render_cover(prs, data, images_by_id: Dict[str, ImageAsset]):
    slide = _blank_slide(prs)
    img = images_by_id.get(data.get("image_id"))
    if img:
        add_picture_contain(slide, img.path, 0, 0, SLIDE_WIDTH, Mm(190), frame=False)
        add_rect(slide, 0, Mm(178), SLIDE_WIDTH, Mm(2), COLOR_GOLD)
    add_rect(slide, 0, Mm(190), SLIDE_WIDTH, SLIDE_HEIGHT - Mm(190), COLOR_NAVY)

    add_text(slide, MARGIN, Mm(198), CONTENT_W, Mm(28), data["title"],
              size=FONT_SIZE_COVER_TITLE, color=COLOR_WHITE, bold=True)
    add_text(slide, MARGIN, Mm(228), CONTENT_W, Mm(10), data["subtitle"],
              size=FONT_SIZE_COVER_SUB, color=COLOR_GRAY_LIGHT)

    highlights = data.get("highlights", [])
    n = max(len(highlights), 1)
    box_w = CONTENT_W / n
    for i, h in enumerate(highlights):
        cx = MARGIN + Emu(int(box_w) * i)
        add_rect(slide, cx, Mm(250), Emu(int(box_w) - Mm(3)), Mm(28), COLOR_NAVY_LIGHT)
        add_text(slide, cx, Mm(254), Emu(int(box_w) - Mm(3)), Mm(20), h,
                  size=Pt(11.5), color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.MIDDLE)
    return slide


def render_toc(prs, data, images_by_id, page_no):
    slide = _blank_slide(prs)
    _header(slide, "목차", page_no)
    y = Mm(35)
    for i, item in enumerate(data["items"], start=1):
        add_rect(slide, MARGIN, y, Mm(9), Mm(9), COLOR_GOLD)
        add_text(slide, MARGIN + Mm(2), y, Mm(6), Mm(9), str(i), size=Pt(12), color=COLOR_WHITE,
                  bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, MARGIN + Mm(14), y, CONTENT_W - Mm(14), Mm(9), item,
                  size=Pt(15), color=COLOR_TEXT_BODY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, MARGIN, y + Mm(13), CONTENT_W, Pt(0.75), COLOR_GRAY_LIGHT)
        y += Mm(20)
    return slide


def _photo_grid(slide, image_ids, images_by_id, top, area_h, captions=None, cols=2):
    ids = [i for i in image_ids if i in images_by_id]
    if not ids:
        return
    rows = (len(ids) + cols - 1) // cols
    gap = Mm(4)
    cell_w = (CONTENT_W - gap * (cols - 1)) / cols
    cell_h = (area_h - gap * (rows - 1)) / rows
    for idx, iid in enumerate(ids):
        r, c = divmod(idx, cols)
        x = MARGIN + Emu(int(cell_w + gap) * c)
        y = top + Emu(int(cell_h + gap) * r)
        img = images_by_id[iid]
        _, box = add_picture_contain(slide, img.path, x, y, int(cell_w), int(cell_h) - Mm(7))
        cap = captions[idx] if captions and idx < len(captions) else _neutral_caption(idx)
        add_text(slide, x, y + int(cell_h) - Mm(6), int(cell_w), Mm(6), cap,
                  size=FONT_SIZE_CAPTION, color=COLOR_TEXT_SUB, align=PP_ALIGN.CENTER)


def render_reasons(prs, data, images_by_id, page_no):
    slide = _blank_slide(prs)
    _header(slide, data["title"], page_no)
    labels = []
    for iid in data["image_ids"]:
        img = images_by_id.get(iid)
        cat = img.category if img else ""
        label_map = {
            "외벽_하자_균열": "외벽 균열", "외벽_하자_박리": "박리 및 들뜸",
            "외벽_하자_오염": "오염 및 얼룩", "외벽_하자_곰팡이": "곰팡이 및 백태",
            "외벽_하자_변색": "변색 및 백화",
        }
        labels.append(label_map.get(cat, "외벽 손상"))
    _photo_grid(slide, data["image_ids"], images_by_id, Mm(28), Mm(220), captions=labels, cols=2)
    add_rect(slide, MARGIN, Mm(258), CONTENT_W, Mm(24), COLOR_GRAY_LIGHT)
    add_text(slide, MARGIN + Mm(4), Mm(263), CONTENT_W - Mm(8), Mm(16), data["note"],
              size=FONT_SIZE_BODY, color=COLOR_TEXT_BODY, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def render_method_overview(prs, data, images_by_id, page_no):
    slide = _blank_slide(prs)
    _header(slide, data["title"], page_no)
    items = data["items"]
    imgs = [i for i in data["image_ids"] if i in images_by_id]
    y = Mm(28)
    row_h = Mm(45)
    for idx, item in enumerate(items):
        add_rect(slide, MARGIN, y, Mm(3), row_h - Mm(4), COLOR_GOLD)
        img_id = imgs[idx] if idx < len(imgs) else None
        if img_id:
            _, box = add_picture_contain(slide, images_by_id[img_id].path,
                                          MARGIN + Mm(8), y, Mm(38), row_h - Mm(4))
            text_x = MARGIN + Mm(50)
        else:
            text_x = MARGIN + Mm(8)
        add_text(slide, text_x, y + Mm(6), CONTENT_W - (text_x - MARGIN) - Mm(2), row_h - Mm(10),
                  item, size=Pt(14), color=COLOR_TEXT_BODY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        y += row_h
    return slide


def render_features(prs, data, images_by_id, page_no):
    slide = _blank_slide(prs)
    _header(slide, data["title"], page_no)
    items = data["items"]
    imgs = [i for i in data["image_ids"] if i in images_by_id]
    y = Mm(28)
    row_h = Mm(47)
    for idx, (title, desc) in enumerate(items):
        img_id = imgs[idx] if idx < len(imgs) else None
        if img_id:
            _, box = add_picture_contain(slide, images_by_id[img_id].path,
                                          MARGIN, y, Mm(42), row_h - Mm(5))
            text_x = MARGIN + Mm(46)
        else:
            add_rect(slide, MARGIN, y, Mm(42), row_h - Mm(5), COLOR_GRAY_LIGHT)
            text_x = MARGIN + Mm(46)
        add_text(slide, text_x, y + Mm(2), CONTENT_W - Mm(46), Mm(8), title,
                  size=Pt(14.5), color=COLOR_NAVY, bold=True)
        add_text(slide, text_x, y + Mm(11), CONTENT_W - Mm(46), row_h - Mm(15), desc,
                  size=FONT_SIZE_BODY, color=COLOR_TEXT_BODY)
        y += row_h
    return slide


def render_process(prs, data, images_by_id, page_no):
    slide = _blank_slide(prs)
    _header(slide, data["title"], page_no)
    steps = data["steps"]
    y = Mm(27)
    row_h = min(Mm(34), Emu(int((SLIDE_HEIGHT - Mm(35)) / max(len(steps), 1))))
    for idx, step in enumerate(steps, start=1):
        add_rect(slide, MARGIN, y, Mm(9), Mm(9), COLOR_NAVY)
        add_text(slide, MARGIN, y, Mm(9), Mm(9), str(idx), size=Pt(12), color=COLOR_WHITE,
                  bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        name = step.name + ("  (참고: 자료 간 표기 상이)" if step.needs_user_confirmation else "")
        text_x = MARGIN + Mm(13)
        img_id = step.image_ids[0] if step.image_ids else None
        img_w = Mm(28)
        if img_id and img_id in images_by_id:
            add_picture_contain(slide, images_by_id[img_id].path,
                                 MARGIN + CONTENT_W - img_w, y, img_w, row_h - Mm(3))
            text_w = CONTENT_W - Mm(13) - img_w - Mm(4)
        else:
            text_w = CONTENT_W - Mm(13)
        add_text(slide, text_x, y, text_w, Mm(7), name, size=Pt(13), color=COLOR_NAVY, bold=True)
        add_text(slide, text_x, y + Mm(8), text_w, row_h - Mm(11), step.description,
                  size=FONT_SIZE_CAPTION, color=COLOR_TEXT_SUB)
        add_rect(slide, MARGIN, y + row_h - Mm(2), CONTENT_W, Pt(0.75), COLOR_GRAY_LIGHT)
        y += row_h
    return slide


def render_case(prs, data, images_by_id, page_no):
    slide = _blank_slide(prs)
    _header(slide, data["title"], page_no)
    add_text(slide, MARGIN, Mm(24), CONTENT_W, Mm(7), data["work_desc"],
              size=Pt(13), color=COLOR_GOLD, bold=True)

    pair = data["pair"]
    before = images_by_id.get(pair.before_image_id)
    after = images_by_id.get(pair.after_image_id)
    half_w = (CONTENT_W - Mm(4)) / 2
    top = Mm(33)
    ph = Mm(85)
    if before:
        add_picture_contain(slide, before.path, MARGIN, top, int(half_w), ph)
        add_rect(slide, MARGIN, top + ph + Mm(2), Emu(int(half_w)), Mm(8), COLOR_NAVY)
        add_text(slide, MARGIN, top + ph + Mm(2), Emu(int(half_w)), Mm(8), "시공 전",
                  size=Pt(12), color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if after:
        ax = MARGIN + Emu(int(half_w + Mm(4)))
        add_picture_contain(slide, after.path, ax, top, int(half_w), ph)
        add_rect(slide, ax, top + ph + Mm(2), Emu(int(half_w)), Mm(8), COLOR_GOLD)
        add_text(slide, ax, top + ph + Mm(2), Emu(int(half_w)), Mm(8), "시공 후",
                  size=Pt(12), color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    proc_top = top + ph + Mm(14)
    if pair.process_image_ids:
        _photo_grid(slide, pair.process_image_ids, images_by_id, proc_top, Mm(60),
                    captions=["공정별 시공 예시"] * len(pair.process_image_ids), cols=3)
        note_top = proc_top + Mm(64)
    else:
        note_top = proc_top

    add_rect(slide, MARGIN, note_top, CONTENT_W, Mm(18), COLOR_GRAY_LIGHT)
    add_text(slide, MARGIN + Mm(4), note_top + Mm(3), CONTENT_W - Mm(8), Mm(12), data["note"],
              size=FONT_SIZE_BODY, color=COLOR_TEXT_BODY, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def render_effects(prs, data, images_by_id, page_no):
    slide = _blank_slide(prs)
    _header(slide, data["title"], page_no)
    _photo_grid(slide, data["image_ids"], images_by_id, Mm(28), Mm(110),
                captions=[_neutral_caption(i) for i in range(len(data["image_ids"]))], cols=2)
    y = Mm(148)
    items = data["items"]
    cols = 2
    cell_w = CONTENT_W / cols
    for idx, item in enumerate(items):
        r, c = divmod(idx, cols)
        x = MARGIN + Emu(int(cell_w) * c)
        yy = y + Emu(int(Mm(18)) * r)
        add_rect(slide, x, yy, Mm(4), Mm(4), COLOR_GOLD)
        add_text(slide, x + Mm(7), yy - Mm(2), int(cell_w) - Mm(9), Mm(10), item,
                  size=Pt(13.5), color=COLOR_TEXT_BODY, bold=True)
    return slide


def render_closing(prs, data, page_no):
    slide = _blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, COLOR_NAVY)
    add_rect(slide, MARGIN, Mm(140), CONTENT_W, Pt(2), COLOR_GOLD)
    add_text(slide, MARGIN, Mm(150), CONTENT_W, Mm(40), data["message"],
              size=Pt(20), color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    return slide


RENDERERS = {
    "reasons": render_reasons,
    "method_overview": render_method_overview,
    "features": render_features,
    "process": render_process,
    "case": render_case,
    "effects": render_effects,
}


def generate_pptx(slide_plan: List[Dict], images_by_id: Dict[str, ImageAsset], out_path: str) -> str:
    prs = new_presentation()
    page_no = 0
    for item in slide_plan:
        t = item["type"]
        if t == "cover":
            render_cover(prs, item, images_by_id)
            continue
        page_no += 1
        if t == "toc":
            render_toc(prs, item, images_by_id, page_no)
        elif t == "closing":
            render_closing(prs, item, page_no)
        elif t in RENDERERS:
            RENDERERS[t](prs, item, images_by_id, page_no)

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return out_path
